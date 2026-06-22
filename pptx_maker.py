"""
기계설비 성능점검 검토자문서 PPTX 생성 모듈 (Pretendard 폰트 내장)
"""

import io
import os
import uuid
import zipfile
import requests
from datetime import datetime
from lxml import etree

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ── 색상 팔레트 ───────────────────────────────────────────────
C_NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
C_BLUE   = RGBColor(0x2D, 0x9C, 0xDB)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_RED    = RGBColor(0xEB, 0x57, 0x57)
C_GRAY   = RGBColor(0x6B, 0x7A, 0x99)
C_LGRAY  = RGBColor(0xF4, 0xF6, 0xF9)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x1A, 0x1A, 0x2E)

FONT_NAME = "Pretendard"
FONT_DIR  = os.path.join(os.path.dirname(__file__), "fonts", "pretendard")

_PRETENDARD_ZIP = (
    "https://github.com/orioncactus/pretendard/releases/download/"
    "v1.3.9/Pretendard-1.3.9.zip"
)
_FONT_FILES = {
    "public/static/alternative/Pretendard-Regular.ttf":  "Pretendard-Regular.ttf",
    "public/static/alternative/Pretendard-Bold.ttf":     "Pretendard-Bold.ttf",
    "public/static/alternative/Pretendard-SemiBold.ttf": "Pretendard-SemiBold.ttf",
    "public/static/alternative/Pretendard-Light.ttf":    "Pretendard-Light.ttf",
}


def _ensure_fonts() -> bool:
    """폰트 파일이 없으면 GitHub 릴리즈에서 자동 다운로드. 성공 여부 반환."""
    required = [os.path.join(FONT_DIR, n) for n in ("Pretendard-Regular.ttf", "Pretendard-Bold.ttf")]
    if all(os.path.exists(p) for p in required):
        return True
    try:
        os.makedirs(FONT_DIR, exist_ok=True)
        resp = requests.get(_PRETENDARD_ZIP, timeout=30)
        resp.raise_for_status()
        zf = zipfile.ZipFile(io.BytesIO(resp.content))
        for src, dst_name in _FONT_FILES.items():
            dst = os.path.join(FONT_DIR, dst_name)
            with open(dst, "wb") as f:
                f.write(zf.read(src))
        return True
    except Exception:
        return False

# ── 슬라이드 치수 (16:9 와이드) ───────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ══════════════════════════════════════════════════════════════
# 폰트 내장 (OOXML ECMA-376 obfuscation)
# ══════════════════════════════════════════════════════════════

def _obfuscate_font(font_bytes: bytes, guid: str) -> bytes:
    """OOXML 표준(ECMA-376 §17.17.9) 폰트 난독화."""
    guid_clean = guid.replace("{", "").replace("}", "").replace("-", "")
    key = bytes.fromhex(guid_clean)
    data = bytearray(font_bytes)
    for i in range(min(32, len(data))):
        data[i] ^= key[15 - (i % 16)]
    return bytes(data)


def _embed_fonts(pptx_bytes: bytes) -> bytes:
    """Pretendard Regular·Bold를 PPTX 패키지에 내장합니다."""
    regular_path = os.path.join(FONT_DIR, "Pretendard-Regular.ttf")
    bold_path    = os.path.join(FONT_DIR, "Pretendard-Bold.ttf")

    if not (os.path.exists(regular_path) and os.path.exists(bold_path)):
        return pptx_bytes   # 폰트 없으면 내장 생략

    guid_r = str(uuid.uuid4()).upper()
    guid_b = str(uuid.uuid4()).upper()

    with open(regular_path, "rb") as f:
        reg_data = _obfuscate_font(f.read(), guid_r)
    with open(bold_path, "rb") as f:
        bold_data = _obfuscate_font(f.read(), guid_b)

    RELS_TYPE_FONT = (
        "http://schemas.openxmlformats.org/officeDocument/2006/"
        "relationships/font"
    )
    CT_FONT = (
        "application/vnd.openxmlformats-officedocument."
        "obfuscatedFont"
    )

    src = io.BytesIO(pptx_bytes)
    dst = io.BytesIO()

    with zipfile.ZipFile(src, "r") as zin, \
         zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:

        existing = {i.filename for i in zin.infolist()}
        rid_r, rid_b = "rIdFontR", "rIdFontB"

        for item in zin.infolist():
            data = zin.read(item.filename)

            # ── presentation.xml: embeddedFontLst 추가 ─────────
            if item.filename == "ppt/presentation.xml":
                root = etree.fromstring(data)
                nsmap = {
                    "p":  "http://schemas.openxmlformats.org/presentationml/2006/main",
                    "a":  "http://schemas.openxmlformats.org/drawingml/2006/main",
                    "r":  "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
                }
                pPr = root.find(qn("p:presentationPr"))
                if pPr is None:
                    pPr = etree.SubElement(root, qn("p:presentationPr"))

                lst = etree.SubElement(pPr, qn("p:embeddedFontLst"))
                ef = etree.SubElement(lst, qn("p:embeddedFont"))
                font_el = etree.SubElement(ef, qn("p:font"))
                font_el.set("typeface", FONT_NAME)
                font_el.set("charset", "0")

                reg_el = etree.SubElement(ef, qn("p:regular"))
                reg_el.set(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id",
                    rid_r,
                )
                bld_el = etree.SubElement(ef, qn("p:bold"))
                bld_el.set(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id",
                    rid_b,
                )
                data = etree.tostring(
                    root, xml_declaration=True, encoding="UTF-8", standalone=True
                )

            # ── presentation.xml.rels: 폰트 관계 추가 ──────────
            elif item.filename == "ppt/_rels/presentation.xml.rels":
                root = etree.fromstring(data)
                for rid, fname in [(rid_r, "font1.fntdata"), (rid_b, "font2.fntdata")]:
                    rel = etree.SubElement(root, "Relationship")
                    rel.set("Id", rid)
                    rel.set("Type", RELS_TYPE_FONT)
                    rel.set("Target", f"fonts/{fname}")
                data = etree.tostring(
                    root, xml_declaration=True, encoding="UTF-8", standalone=True
                )

            # ── [Content_Types].xml: font MIME 추가 ────────────
            elif item.filename == "[Content_Types].xml":
                root = etree.fromstring(data)
                CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
                existing_parts = {
                    el.get("PartName")
                    for el in root.findall(f"{{{CT_NS}}}Override")
                }
                for fname in ("font1.fntdata", "font2.fntdata"):
                    part_name = f"/ppt/fonts/{fname}"
                    if part_name not in existing_parts:
                        ov = etree.SubElement(root, f"{{{CT_NS}}}Override")
                        ov.set("PartName", part_name)
                        ov.set("ContentType", CT_FONT)
                data = etree.tostring(
                    root, xml_declaration=True, encoding="UTF-8", standalone=True
                )

            zout.writestr(item, data)

        zout.writestr("ppt/fonts/font1.fntdata", reg_data)
        zout.writestr("ppt/fonts/font2.fntdata", bold_data)

    dst.seek(0)
    return dst.read()


# ══════════════════════════════════════════════════════════════
# 유틸: 텍스트 박스 / 도형 헬퍼
# ══════════════════════════════════════════════════════════════

def _set_font(run, size: int, bold=False, color: RGBColor = C_BLACK):
    run.font.name   = FONT_NAME
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.color.rgb = color


def _add_textbox(slide, left, top, width, height,
                 text="", size=12, bold=False,
                 color=C_BLACK, align=PP_ALIGN.LEFT,
                 word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = word_wrap
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size, bold, color)
    return txb


def _add_rect(slide, left, top, width, height,
              fill: RGBColor = None, line: RGBColor = None, line_w: int = 0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()

    if line and line_w:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _fill_rect_text(slide, left, top, width, height,
                    text, size=11, bold=False,
                    bg: RGBColor = C_LGRAY, fg: RGBColor = C_BLACK,
                    align=PP_ALIGN.LEFT, padding_left=Inches(0.12)):
    _add_rect(slide, left, top, width, height, fill=bg)
    _add_textbox(slide, left + padding_left, top, width - padding_left,
                 height, text, size, bold, fg, align)


# ══════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ══════════════════════════════════════════════════════════════

def _slide_cover(prs: Presentation, data: dict, today: str):
    blank = prs.slide_layouts[6]   # 완전 빈 레이아웃
    slide = prs.slides.add_slide(blank)

    # 전체 배경 (네이비)
    _add_rect(slide, 0, 0, W, H, fill=C_NAVY)

    # 상단 수평선 장식
    _add_rect(slide, 0, 0, W, Inches(0.06), fill=C_BLUE)

    # 메인 타이틀
    _add_textbox(
        slide,
        Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.6),
        "기계설비 성능점검\n검토자문서",
        size=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
    )

    # 구분선
    _add_rect(slide, Inches(4.5), Inches(3.7), Inches(4.33), Inches(0.04), fill=C_BLUE)

    basic = data.get("basic_info", {})
    building = basic.get("building_name", "")
    address  = basic.get("building_address", "")

    _add_textbox(
        slide,
        Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.7),
        building, size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(1.5), Inches(4.65), Inches(10.3), Inches(0.5),
        address, size=14, bold=False, color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER,
    )

    # 날짜
    _add_textbox(
        slide,
        Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.5),
        today, size=13, bold=False, color=RGBColor(0x90, 0xA8, 0xC8), align=PP_ALIGN.CENTER,
    )

    # 하단 장식
    _add_rect(slide, 0, Inches(7.3), W, Inches(0.2), fill=C_BLUE)


# ══════════════════════════════════════════════════════════════
# 슬라이드 2: 기본정보 + 요약
# ══════════════════════════════════════════════════════════════

def _slide_summary(prs: Presentation, data: dict, items: list):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 헤더 바
    _add_rect(slide, 0, 0, W, Inches(1.0), fill=C_NAVY)
    _add_textbox(
        slide,
        Inches(0.4), Inches(0.22), Inches(8), Inches(0.6),
        "검토 요약", size=22, bold=True, color=C_WHITE,
    )

    basic = data.get("basic_info", {})
    adequate      = sum(1 for i in items if i.get("ai_judgment") == "적정")
    needs_improve = len(items) - adequate

    # ── 좌측: 기본정보 테이블 ─────────────────────────────────
    rows = [
        ("건축물명",     basic.get("building_name",         "확인 불가")),
        ("주소",         basic.get("building_address",       "확인 불가")),
        ("관리주체",     basic.get("management_entity",      "확인 불가")),
        ("점검업체",     basic.get("inspection_company",     "확인 불가")),
        ("업체 대표자",  basic.get("company_representative", "확인 불가")),
    ]

    row_h = Inches(0.58)
    col_l_w = Inches(1.55)
    col_v_w = Inches(5.9)
    start_y = Inches(1.15)

    for idx, (label, val) in enumerate(rows):
        y = start_y + idx * row_h
        bg = C_LGRAY if idx % 2 == 0 else C_WHITE
        _fill_rect_text(
            slide, Inches(0.3), y, col_l_w, row_h,
            label, size=10, bold=True, bg=C_NAVY, fg=C_WHITE,
        )
        _fill_rect_text(
            slide, Inches(0.3) + col_l_w, y, col_v_w, row_h,
            val, size=10, bg=bg, fg=C_BLACK,
        )

    # ── 우측: 통계 카드 ──────────────────────────────────────
    card_y  = Inches(1.3)
    card_h  = Inches(1.4)
    card_x1 = Inches(8.2)
    card_x2 = Inches(10.75)
    card_w  = Inches(2.1)

    def stat_card(slide, x, y, w, h, label, count, total, bg):
        _add_rect(slide, x, y, w, h, fill=bg)
        pct = f"{count/total*100:.0f}%" if total else "0%"
        _add_textbox(slide, x, y + Inches(0.15), w, Inches(0.35),
                     label, size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(slide, x, y + Inches(0.45), w, Inches(0.55),
                     str(count), size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(slide, x, y + Inches(0.95), w, Inches(0.3),
                     pct, size=13, color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)

    # 전체
    _add_rect(slide, Inches(8.2), card_y, Inches(5.0), card_h, fill=C_BLUE)
    _add_textbox(slide, Inches(8.2), card_y + Inches(0.15), Inches(5.0), Inches(0.35),
                 "전체 점검 항목", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(8.2), card_y + Inches(0.45), Inches(5.0), Inches(0.55),
                 str(len(items)), size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(8.2), card_y + Inches(0.95), Inches(5.0), Inches(0.3),
                 "건", size=13, color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)

    # 적정 / 보완필요
    stat_card(slide, card_x1, card_y + Inches(1.6), card_w, card_h,
              "적정", adequate, len(items), C_GREEN)
    stat_card(slide, card_x2, card_y + Inches(1.6), card_w, card_h,
              "보완필요", needs_improve, len(items), C_RED)

    # 막대 비율 표시
    bar_y   = Inches(4.5)
    bar_h   = Inches(0.45)
    bar_x   = Inches(8.2)
    bar_w   = Inches(5.0)
    if len(items) > 0:
        g_w = bar_w * adequate / len(items)
        r_w = bar_w - g_w
        if g_w > 0:
            _add_rect(slide, bar_x, bar_y, g_w, bar_h, fill=C_GREEN)
        if r_w > 0:
            _add_rect(slide, bar_x + g_w, bar_y, r_w, bar_h, fill=C_RED)

    # 하단 선
    _add_rect(slide, 0, Inches(7.3), W, Inches(0.05), fill=C_BLUE)


# ══════════════════════════════════════════════════════════════
# 슬라이드 N: 점검 항목 상세 (2개씩 묶음)
# ══════════════════════════════════════════════════════════════

def _item_block(slide, item: dict, top: float, block_h: float):
    """한 항목 블록을 슬라이드에 그린다 (top 기준 block_h 높이)."""
    pad   = Inches(0.25)
    left  = Inches(0.3)
    right = W - Inches(0.3)
    width = right - left

    judgment = item.get("ai_judgment", "")
    is_ok    = judgment == "적정"
    accent   = C_GREEN if is_ok else C_RED

    # 항목 헤더
    header_h = Inches(0.42)
    _add_rect(slide, left, top, width, header_h, fill=accent)
    label = f"[{judgment}]  [{item.get('page_number','')}] {item.get('item_name','')}"
    _add_textbox(slide, left + pad, top, width - pad, header_h,
                 label, size=12, bold=True, color=C_WHITE)

    # 필드 그리드
    field_top  = top + header_h
    field_h    = Inches(0.38)
    col1_w     = Inches(1.6)
    col2_w     = (width / 2) - col1_w
    col3_w     = Inches(1.6)
    col4_w     = (width / 2) - col3_w

    fields_left = [
        ("점검 기준", item.get("criteria", "")),
        ("업체 점검결과", item.get("reported_result", "")),
    ]
    fields_right = [
        ("사진 판독값", item.get("photo_value", "")),
        ("보고서 기재값", item.get("report_value", "")),
    ]

    for r_idx, ((l_lbl, l_val), (r_lbl, r_val)) in enumerate(
        zip(fields_left, fields_right)
    ):
        y   = field_top + r_idx * field_h
        bg  = C_LGRAY if r_idx % 2 == 0 else C_WHITE
        mid = left + width / 2

        _fill_rect_text(slide, left,          y, col1_w, field_h, l_lbl, 9, True,  C_NAVY, C_WHITE)
        _fill_rect_text(slide, left + col1_w, y, col2_w, field_h, l_val, 9, False, bg,     C_BLACK)
        _fill_rect_text(slide, mid,            y, col3_w, field_h, r_lbl, 9, True,  C_NAVY, C_WHITE)
        _fill_rect_text(slide, mid + col3_w,  y, col4_w, field_h, r_val, 9, False, bg,     C_BLACK)

    # 검토자문의견
    opinion_top = field_top + 2 * field_h
    opinion_h   = block_h - header_h - 2 * field_h - Inches(0.05)

    _fill_rect_text(slide, left,            opinion_top, col1_w, opinion_h,
                    "검토자문의견", 9, True, C_NAVY, C_WHITE)
    opinion_bg = RGBColor(0xF0, 0xFB, 0xF4) if is_ok else RGBColor(0xFF, 0xF2, 0xF2)
    _fill_rect_text(slide, left + col1_w, opinion_top, width - col1_w, opinion_h,
                    item.get("expert_detailed_opinion", ""),
                    9, False, opinion_bg, C_BLACK)


def _slides_items(prs: Presentation, items: list, title="점검 항목"):
    target = items

    blank   = prs.slide_layouts[6]
    hdr_h   = Inches(0.85)
    per_pg  = 2
    block_h = (H - hdr_h - Inches(0.1)) / per_pg

    for page_idx in range(0, max(1, len(target)), per_pg):
        chunk = target[page_idx: page_idx + per_pg]
        slide = prs.slides.add_slide(blank)

        # 헤더
        _add_rect(slide, 0, 0, W, hdr_h, fill=C_NAVY)
        start = page_idx + 1
        end   = min(page_idx + per_pg, len(target))
        _add_textbox(
            slide,
            Inches(0.4), Inches(0.18), Inches(10), Inches(0.55),
            f"{title}  ({start}–{end} / 총 {len(target)}건)",
            size=18, bold=True, color=C_WHITE,
        )
        _add_rect(slide, 0, Inches(7.3), W, Inches(0.05), fill=C_BLUE)

        for i, item in enumerate(chunk):
            top = hdr_h + i * block_h
            _item_block(slide, item, top, block_h)

        if not chunk:
            _add_textbox(
                slide,
                Inches(2), Inches(3), Inches(9), Inches(1),
                "해당 항목이 없습니다.", size=18, color=C_GRAY, align=PP_ALIGN.CENTER,
            )


# ══════════════════════════════════════════════════════════════
# 공개 API
# ══════════════════════════════════════════════════════════════

def create_review_pptx(data: dict) -> bytes:
    """
    data: {"basic_info": {...}, "items": [...]}
    반환: PPTX 파일 바이트 (Pretendard 폰트 내장)
    """
    _ensure_fonts()

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    items = data.get("items", [])
    today = datetime.now().strftime("%Y년 %m월 %d일")

    adequate_items   = [i for i in items if i.get("ai_judgment") == "적정"]
    inadequate_items = [i for i in items if i.get("ai_judgment") != "적정"]

    _slide_cover(prs, data, today)
    _slide_summary(prs, data, items)
    _slides_items(prs, adequate_items,   title="🟢 적정 항목")
    _slides_items(prs, inadequate_items, title="🔴 보완필요 항목")

    buf = io.BytesIO()
    prs.save(buf)
    result = _embed_fonts(buf.getvalue())
    return result
