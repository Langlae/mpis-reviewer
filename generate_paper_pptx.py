#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
학술발표대회 발표자료 생성 스크립트
AI 기반 기계설비 성능점검 결과보고서 자동검토 시스템 개발
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# ── 색상 상수 ────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1B, 0x3A, 0x6B)
BLUE  = RGBColor(0x2D, 0x9C, 0xDB)
LBLUE = RGBColor(0xEB, 0xF5, 0xFB)
GRN   = RGBColor(0x27, 0xAE, 0x60)
LGRN  = RGBColor(0xE9, 0xF7, 0xEF)
RED   = RGBColor(0xEB, 0x57, 0x57)
LRED  = RGBColor(0xFD, 0xED, 0xED)
ORG   = RGBColor(0xF2, 0x99, 0x4A)
LORG  = RGBColor(0xFE, 0xF9, 0xE7)
PRP   = RGBColor(0x8E, 0x44, 0xAD)
LGRAY = RGBColor(0xF5, 0xF6, 0xF7)
MGRAY = RGBColor(0xDD, 0xDD, 0xDD)
DGRAY = RGBColor(0x44, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── 헬퍼 함수 ─────────────────────────────────────────────────────────
def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def r(slide, l, t, w, h, fill=None, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def t(slide, l, top, w, h, text, sz=11, bold=False, color=None,
      align=PP_ALIGN.LEFT, fn='맑은 고딕', wrap=True):
    tb = slide.shapes.add_textbox(Cm(l), Cm(top), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.name = fn
        if color:
            run.font.color.rgb = color
    return tb

def ph(slide, l, top, w, h, label='이미지 삽입', sz=10):
    """이미지 자리표시자 (점선 테두리)"""
    s = r(slide, l, top, w, h, LGRAY, MGRAY, Pt(1.5))
    ln = s.line._ln
    for old in ln.findall(qn('a:prstDash')):
        ln.remove(old)
    pd = etree.SubElement(ln, qn('a:prstDash'))
    pd.set('val', 'lgDash')
    t(slide, l, top + h/2 - 0.45, w, 0.9,
      f'[ {label} ]', sz, False, RGBColor(0xAA, 0xAA, 0xAA), PP_ALIGN.CENTER)
    return s

def hdr(slide, title, num, total=8):
    r(slide, 0, 0, 33.87, 2.0, NAVY)
    r(slide, 0, 1.85, 33.87, 0.15, BLUE)
    t(slide, 1.3, 0.28, 26.0, 1.5, title, 20, True, WHITE)
    t(slide, 28.0, 0.28, 5.5, 1.5, f'{num} / {total}', 11, False,
      RGBColor(0x88, 0xAA, 0xCC), PP_ALIGN.RIGHT)

# ── 프레젠테이션 생성 ────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
blank = prs.slide_layouts[6]

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
bg(s1, NAVY)

r(s1, 0, 0, 33.87, 0.55, BLUE)
r(s1, 0, 17.6, 33.87, 1.45, RGBColor(0x0F, 0x27, 0x4A))

# 왼쪽 세로 강조선
r(s1, 0, 0.55, 0.55, 17.05, RGBColor(0x12, 0x2A, 0x52))

# 기관 로고 자리표시자
r(s1, 1.3, 1.0, 6.0, 1.5, WHITE)
t(s1, 1.3, 1.1, 6.0, 1.3, '대한기계설비산업연구원\n(KRIMFI)', 9, True, NAVY, PP_ALIGN.CENTER)

# AI 배지
r(s1, 26.5, 1.0, 6.0, 1.5, BLUE)
t(s1, 26.5, 1.1, 6.0, 1.3, 'AI·자동화 기술 연구', 10, True, WHITE, PP_ALIGN.CENTER)

# 구분선
r(s1, 1.5, 4.3, 30.87, 0.1, BLUE)

# 메인 제목
t(s1, 1.5, 4.7, 30.87, 4.2,
  'AI 기반 기계설비 성능점검\n결과보고서 자동검토 시스템 개발',
  30, True, WHITE, PP_ALIGN.CENTER)

# 영문 부제
t(s1, 1.5, 9.2, 30.87, 1.7,
  'Development of an AI-Based Automatic Review System\nfor Mechanical Equipment Performance Inspection Reports',
  14, False, RGBColor(0xA0, 0xB8, 0xD0), PP_ALIGN.CENTER)

# 구분선 2
r(s1, 10.0, 11.2, 13.87, 0.08, BLUE)

# 저자
t(s1, 1.5, 11.5, 30.87, 1.0,
  '[제1저자명]1,  [공동저자명]1,  [교신저자명]2†',
  13, False, WHITE, PP_ALIGN.CENTER)

t(s1, 1.5, 12.5, 30.87, 0.9,
  '1,2대한기계설비산업연구원',
  11, False, RGBColor(0xB0, 0xC4, 0xDE), PP_ALIGN.CENTER)

# 학술대회 정보
t(s1, 1.5, 17.8, 30.87, 0.7,
  '2026년 대한기계설비산업연구원 학술발표대회',
  10, False, RGBColor(0x88, 0xA8, 0xC8), PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — 목차
# ════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
bg(s2, WHITE)
hdr(s2, '목  차', 2)

toc = [
    ('01', '연구 배경 및 필요성', '기계설비 성능점검 제도 현황 및 문제점'),
    ('02', '연구 목표 및 접근 방법', '기존 수동 검토의 한계와 AI 해결 방안'),
    ('03', '시스템 설계', '멀티모달 AI 기반 자동검토 아키텍처'),
    ('04', 'AI 검토 5대 원칙', '정량적 검토의견 생성을 위한 원칙 설계'),
    ('05', '시스템 구현', '웹 기반 플랫폼 구현 및 주요 화면'),
    ('06', '결론 및 향후 연구', '기대효과 및 후속 연구 계획'),
]

for i, (num, title, desc) in enumerate(toc):
    col = i % 2
    row = i // 2
    x = 1.5 + col * 16.2
    y = 2.5 + row * 2.9

    r(s2, x, y, 2.0, 2.0, NAVY)
    t(s2, x, y + 0.35, 2.0, 1.3, num, 20, True, WHITE, PP_ALIGN.CENTER)

    r(s2, x + 2.2, y, 13.7, 2.0, LGRAY)
    t(s2, x + 2.4, y + 0.2, 13.3, 0.85, title, 13, True, NAVY)
    t(s2, x + 2.4, y + 1.1, 13.3, 0.75, desc, 10, False, DGRAY)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — 연구 배경 및 필요성
# ════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
bg(s3, WHITE)
hdr(s3, '연구 배경 및 필요성', 3)

# 제도 현황 박스
r(s3, 1.2, 2.3, 31.0, 1.8, LBLUE)
r(s3, 1.2, 2.3, 0.35, 1.8, BLUE)
t(s3, 1.8, 2.4, 30.2, 1.6,
  '「기계설비법」 제17조: 일정 규모 이상 건축물 관리주체는 매년 1회 이상 성능점검 실시 의무\n'
  '(2021년 8월 시행, 연면적 규모별 적용 대상 순차 확대 → 연간 보고서 수 지속 증가)',
  11, False, DGRAY)

# 문제점 카드 4개
probs = [
    (RED,  '①', '보고서 품질 편차',
     '업체 간 기술역량 차이\n점검항목 해석 상이\n보고서 서식 비표준화'),
    (ORG,  '②', '검토 역량 한계',
     '지자체 공무원\n기술 전문성 부족\n행정 여력 한계'),
    (PRP,  '③', '검토 수요 증가',
     '성능점검 대상 건축물\n지속 증가\n전문위원 업무 가중'),
    (NAVY, '④', '수작업 비효율',
     '수십~수백 항목\n수작업 검토\n주관적 판단 편차'),
]

for i, (color, num, title, desc) in enumerate(probs):
    x = 1.2 + i * 7.9
    y = 4.5
    r(s3, x, y, 7.4, 6.3, LGRAY)
    r(s3, x, y, 7.4, 1.6, color)
    t(s3, x, y + 0.07, 2.2, 1.5, num, 22, True, WHITE, PP_ALIGN.CENTER)
    t(s3, x + 2.2, y + 0.35, 5.0, 0.9, title, 12, True, WHITE)
    t(s3, x + 0.3, y + 1.9, 6.8, 4.1, desc, 11, False, DGRAY)

# 결론 배너
r(s3, 1.2, 11.1, 31.0, 1.7, NAVY)
t(s3, 1.5, 11.2, 30.7, 1.5,
  '→  AI 기반 자동검토 시스템 개발을 통한 검토 효율화 및 보고서 품질 균질화 필요',
  14, True, WHITE, PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — 연구 목표 및 비교
# ════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
bg(s4, WHITE)
hdr(s4, '연구 목표 및 접근 방법', 4)

# 기존 방식
r(s4, 1.0, 2.3, 14.8, 10.5, LGRAY)
r(s4, 1.0, 2.3, 14.8, 1.1, RGBColor(0x78, 0x80, 0x90))
t(s4, 1.0, 2.3, 14.8, 1.1, '기존: 전문가 수동 검토', 13, True, WHITE, PP_ALIGN.CENTER)

manual = [
    ('●', '성능점검 결과보고서 PDF 수동 열람'),
    ('●', '점검항목 수치 육안 확인'),
    ('●', '계측 사진 개별 확인'),
    ('●', '검토의견 수작업 작성'),
    ('●', 'Excel 결과 수동 정리'),
    ('', ''),
    ('⏱', '소요 시간: 수 시간 ~ 수 일'),
    ('⚠', '주관적 판단 편차 존재'),
    ('⚠', '검토량 증가 시 처리 한계'),
]
for i, (icon, item) in enumerate(manual):
    clr = RED if '⚠' in icon else (ORG if '⏱' in icon else DGRAY)
    bld = icon in ('⏱', '⚠')
    t(s4, 1.4, 3.7 + i * 0.88, 14.1, 0.85, f'{icon}  {item}', 11, bld, clr)

# 화살표
t(s4, 16.3, 6.3, 1.4, 2.5, '→\nAI', 16, True, BLUE, PP_ALIGN.CENTER)

# AI 방식
r(s4, 18.0, 2.3, 14.8, 10.5, LBLUE)
r(s4, 18.0, 2.3, 14.8, 1.1, BLUE)
t(s4, 18.0, 2.3, 14.8, 1.1, 'AI 자동검토 시스템', 13, True, WHITE, PP_ALIGN.CENTER)

ai_items = [
    ('●', 'PDF + 이미지 자동 업로드'),
    ('●', '멀티모달 AI 전체 분석'),
    ('●', '계측 사진 수치 자동 추출'),
    ('●', '정량적 검토의견 자동 생성'),
    ('●', 'Excel / PPTX 자동 출력'),
    ('', ''),
    ('⚡', '소요 시간: 수 분 이내'),
    ('✓', '일관된 정량적 판단 기준'),
    ('✓', '대용량 보고서 처리 가능'),
]
for i, (icon, item) in enumerate(ai_items):
    clr = GRN if '✓' in icon else (BLUE if '⚡' in icon else DGRAY)
    bld = icon in ('⚡', '✓')
    t(s4, 18.4, 3.7 + i * 0.88, 14.1, 0.85, f'{icon}  {item}', 11, bld, clr)

# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — 시스템 아키텍처
# ════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
bg(s5, WHITE)
hdr(s5, '시스템 설계 및 아키텍처', 5)

# 입력 박스
r(s5, 0.8, 2.8, 7.3, 9.0, LBLUE)
r(s5, 0.8, 2.8, 7.3, 1.0, BLUE)
t(s5, 0.8, 2.8, 7.3, 1.0, '입  력', 13, True, WHITE, PP_ALIGN.CENTER)
t(s5, 1.0, 4.1, 6.9, 1.8,
  '📄  성능점검 결과보고서\n        (PDF 형식)', 11, False, DGRAY)
r(s5, 1.0, 6.2, 6.9, 0.05, MGRAY)
t(s5, 1.0, 6.5, 6.9, 1.8,
  '📷  현장 계측 사진\n        (JPEG / PNG)', 11, False, DGRAY)

# 화살표
t(s5, 8.4, 6.5, 1.4, 1.5, '→', 26, True, NAVY, PP_ALIGN.CENTER)

# AI 엔진 박스
r(s5, 10.0, 2.3, 10.0, 11.0, LGRAY)
r(s5, 10.0, 2.3, 10.0, 1.0, NAVY)
t(s5, 10.0, 2.3, 10.0, 1.0, 'AI 분석 엔진', 13, True, WHITE, PP_ALIGN.CENTER)
t(s5, 10.0, 3.4, 10.0, 0.85, 'Google Gemini (멀티모달 LLM)', 10, True, BLUE, PP_ALIGN.CENTER)

steps = [
    (BLUE,  '①', '항목 전수 추출',   'PDF에서 모든 점검항목\n100% 추출 보장'),
    (GRN,   '②', '교차검증',         '계측 사진 수치 vs\n보고서 기재값 비교'),
    (ORG,   '③', '검토의견 생성',    '5대 원칙 기반\n정량적 의견 작성'),
    (PRP,   '④', '결과 구조화',      'JSON 형식으로\n항목별 데이터 정리'),
]
for i, (color, num, title, desc) in enumerate(steps):
    y = 4.5 + i * 2.2
    r(s5, 10.2, y, 9.6, 1.9, WHITE, color, Pt(0.75))
    r(s5, 10.2, y, 1.5, 1.9, color)
    t(s5, 10.2, y + 0.35, 1.5, 1.2, num, 14, True, WHITE, PP_ALIGN.CENTER)
    t(s5, 11.9, y + 0.1, 7.7, 0.75, title, 12, True, NAVY)
    t(s5, 11.9, y + 0.9, 7.7, 0.9, desc, 10, False, DGRAY)

# 화살표
t(s5, 20.3, 6.5, 1.4, 1.5, '→', 26, True, NAVY, PP_ALIGN.CENTER)

# 출력 박스
r(s5, 21.9, 2.3, 11.0, 11.0, LGRN)
r(s5, 21.9, 2.3, 11.0, 1.0, GRN)
t(s5, 21.9, 2.3, 11.0, 1.0, '출  력', 13, True, WHITE, PP_ALIGN.CENTER)

outputs = [
    ('화면 표시', '• 기본정보 탭 (건물·업체·위치 지도)\n• 적정 항목 탭 (녹색 표시)\n• 보완필요 항목 탭 (적색 표시)'),
    ('Excel 출력', '• 항목별 검토의견 전체\n• 적정/보완필요 구분\n• 설계값/측정값/비율'),
    ('PPTX 출력', '• 요약 통계 슬라이드\n• 항목별 상세 슬라이드\n• 검토자문 보고서 형식'),
]
for i, (title, desc) in enumerate(outputs):
    y = 3.5 + i * 3.2
    r(s5, 22.1, y, 10.6, 2.8, WHITE, GRN, Pt(0.5))
    t(s5, 22.3, y + 0.1, 10.2, 0.75, title, 12, True, GRN)
    t(s5, 22.3, y + 0.9, 10.2, 1.8, desc, 10, False, DGRAY)

# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — AI 검토 5대 원칙
# ════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
bg(s6, WHITE)
hdr(s6, 'AI 검토 5대 원칙', 6)

t(s6, 1.2, 2.2, 31.0, 0.9,
  '정성적 판단("양호", "적합")을 배제하고, 수치 기반의 정량적 전문가 검토의견을 생성하기 위한 원칙을 시스템 프롬프트에 내재화',
  10, False, DGRAY)

principles = [
    (BLUE, '원칙 1', '측정값/설계값 비율 계산',
     '측정값과 설계값의 비율(%)을 명시하여\n판단 근거를 정량화',
     '"측정값 341.74 m³/h = 설계값 370 m³/h의 92.4%\n → ±10% 허용 편차 범위 이내"'),
    (GRN,  '원칙 2', '단위 환산 수행',
     '단위가 상이한 경우 환산 계산\n과정을 명시하여 검증 가능성 확보',
     '"풍속 2.3 m/s × 0.45 m² × 3,600 = 3,726 m³/h\n → 설계값 대비 비율 계산 가능"'),
    (ORG,  '원칙 3', '운전 조건 확인',
     '정격 운전 여부 및 실제 운전 상태를\n구분하여 판단 조건 명시',
     '"정격 주파수 60Hz 미운전 상태 측정치\n → 재측정 요청 필요"'),
    (PRP,  '원칙 4', '정량적 판단 논리 구조',
     '"측정값 X, 설계값 Y, 비율 Z%\n → 기준 충족/미충족" 구조적 서술',
     '"측정값 18,634 CMH, 설계값 21,000 CMH,\n비율 88.7% → ±10% 기준 미충족"'),
    (RED,  '원칙 5', '구체적 보완 방향 제시',
     '단순 "보완 필요" 대신\n구체적 조치 방향 제시',
     '"설계도서에서 설계값 확인 후 성능점검표\n기재 / 정격 상태 재측정 후 결과 갱신"'),
]

# 3개 + 2개 배치
for i, (color, num, title, principle, example) in enumerate(principles):
    if i < 3:
        x = 1.2 + i * 10.8
        y = 3.3
        w = 10.3
    else:
        x = 5.7 + (i - 3) * 11.7
        y = 10.0
        w = 11.2

    r(s6, x, y, w, 5.8, LGRAY)
    r(s6, x, y, w, 1.3, color)
    t(s6, x + 0.2, y + 0.08, 2.8, 1.1, num, 13, True, WHITE)
    t(s6, x + 3.0, y + 0.28, w - 3.2, 0.75, title, 11, True, WHITE)
    t(s6, x + 0.3, y + 1.5, w - 0.5, 1.7, principle, 10, False, DGRAY)
    r(s6, x + 0.2, y + 3.2, w - 0.4, 2.3, RGBColor(0xF0, 0xF4, 0xFF))
    t(s6, x + 0.4, y + 3.3, w - 0.7, 2.1, f'예) {example}', 9, False,
      RGBColor(0x33, 0x33, 0x66), PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — 시스템 구현
# ════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
bg(s7, WHITE)
hdr(s7, '시스템 구현', 7)

t(s7, 1.2, 2.2, 31.0, 0.8,
  'Streamlit 기반 웹 애플리케이션  ·  Google Gemini API (멀티모달)  ·  인터랙티브 재검토 기능',
  10, False, DGRAY, PP_ALIGN.CENTER)

screens = [
    ('파일 업로드 및 AI 분석 화면', 'Fig. 2'),
    ('검토의견 확인 화면\n(적정 / 보완필요 탭)', 'Fig. 3'),
    ('검토자문 보고서 출력\n(Excel / PPTX)', 'Fig. 4'),
]
for i, (label, fignum) in enumerate(screens):
    x = 1.2 + i * 10.9
    ph(s7, x, 3.2, 10.3, 7.8, label, 10)
    r(s7, x, 11.0, 10.3, 0.75, LGRAY)
    t(s7, x, 11.05, 10.3, 0.7, f'{fignum}  {label.split(chr(10))[0]}',
      9, False, DGRAY, PP_ALIGN.CENTER)

# 주요 기능
r(s7, 1.2, 12.1, 31.0, 0.75, NAVY)
t(s7, 1.4, 12.1, 31.0, 0.75, '주요 기능', 12, True, WHITE)
r(s7, 1.2, 12.85, 31.0, 3.0, LGRAY)

feats = [
    '📄  멀티모달 PDF + 이미지 동시 분석 (Google Gemini)',
    '🔄  항목별 인터랙티브 재검토 요청 기능',
    '🔒  JSON 복원 알고리즘으로 100% 항목 추출 보장',
    '🗺  Google Maps 연동 건축물 위치 확인',
    '📊  Excel / PPTX 원클릭 검토자문 보고서 출력',
]
for i, feat in enumerate(feats):
    col = i % 3
    row = i // 3
    t(s7, 1.7 + col * 10.5, 13.1 + row * 1.3, 10.0, 1.2,
      feat, 10, False, DGRAY)

# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — 결론 및 향후 연구
# ════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
bg(s8, WHITE)
hdr(s8, '결론 및 향후 연구', 8)

# 결론 카드 3개
conc = [
    (BLUE, '검토 효율화',
     '기존 수 시간의 수동 검토를\n수 분 이내 자동 완료\n→ 검토자문위원 업무 부담 경감'),
    (GRN,  '품질 균질화',
     '5대 정량적 검토 원칙 내재화\n일관성 있는 수치 기반 판단\n→ 주관적 편차 최소화'),
    (ORG,  '멀티모달 교차검증',
     'PDF + 계측 사진 동시 분석\n보고서 기재값 vs 실측값 검증\n→ 누락·오기 자동 탐지'),
]
for i, (color, title, desc) in enumerate(conc):
    x = 1.2 + i * 10.8
    r(s8, x, 2.3, 10.3, 6.2, LGRAY)
    r(s8, x, 2.3, 10.3, 1.3, color)
    t(s8, x, 2.3, 10.3, 1.3, title, 15, True, WHITE, PP_ALIGN.CENTER)
    t(s8, x + 0.3, 3.8, 9.7, 4.5, desc, 11, False, DGRAY)

# 향후 연구
r(s8, 1.2, 9.0, 31.0, 0.85, NAVY)
t(s8, 1.4, 9.0, 30.7, 0.85, '향후 연구 계획', 13, True, WHITE)

future = [
    (BLUE, '단기',
     '축적된 검토 데이터 기반\nAI 정확도 지속 고도화'),
    (GRN,  '중기',
     '기계설비 도메인 특화\nFine-tuning 모델 적용'),
    (ORG,  '장기',
     'MIS(민원웹포털) 연계\n자동 검토 확인서 생성'),
]
for i, (color, period, desc) in enumerate(future):
    x = 1.2 + i * 10.8
    r(s8, x, 10.1, 10.3, 3.7, LGRAY)
    r(s8, x, 10.1, 2.8, 3.7, color)
    t(s8, x, 10.1, 2.8, 3.7, period, 14, True, WHITE, PP_ALIGN.CENTER)
    t(s8, x + 3.0, 10.5, 7.1, 2.9, desc, 11, False, DGRAY)

# 마무리 배너
r(s8, 1.2, 14.2, 31.0, 2.2, NAVY)
t(s8, 1.2, 14.2, 31.0, 2.2,
  '감사합니다\n본 연구는 대한기계설비산업연구원의 지원으로 수행되었습니다.',
  16, True, WHITE, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════════════
output = r'c:\Users\USER\Desktop\py(성능점검)\mpis-reviewer\학술발표_발표자료.pptx'
prs.save(output)
print(f'저장 완료: {output}')
