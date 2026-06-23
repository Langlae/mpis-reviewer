"""
기계설비 성능점검 검토자문 시스템 개발 — 학술발표 PPTX 생성 스크립트
출력: 검토자문시스템_발표자료.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── 색상 팔레트 ──────────────────────────────────
NAVY   = RGBColor(0x14, 0x30, 0x5A)
NAVY2  = RGBColor(0x1A, 0x3D, 0x70)
NAVY3  = RGBColor(0x22, 0x50, 0x8A)
BLUE   = RGBColor(0x18, 0x68, 0xC8)
GREEN  = RGBColor(0x00, 0x7D, 0x5E)
ORANGE = RGBColor(0xD6, 0x50, 0x1C)
PURPLE = RGBColor(0x60, 0x40, 0xB8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFB)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
BORD   = RGBColor(0xD0, 0xD8, 0xE8)
TEXT   = RGBColor(0x1C, 0x2E, 0x40)
TEXT2  = RGBColor(0x4A, 0x60, 0x70)
TEXT3  = RGBColor(0x8A, 0x9D, 0xAC)
BG3    = RGBColor(0xE8, 0xEF, 0xF7)
BLUEL  = RGBColor(0xE4, 0xEE, 0xF8)
GREENL = RGBColor(0xE4, 0xF5, 0xEE)

W, H = Inches(13.33), Inches(7.5)   # 1280×720 비율

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # blank layout


# ── 헬퍼 함수 ────────────────────────────────────
def add_slide():
    return prs.slides.add_slide(blank)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def rect(slide, x, y, w, h, fill=None, line=None, line_w=None, radius=None):
    """사각형 도형 추가 (inches 단위)"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = Pt(line_w)
        else:
            shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, x, y, w, h, text, size, bold=False, color=None, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    """텍스트 박스 추가 (inches 단위)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb

def hdr(slide, title, slide_num):
    """공통 헤더 (navy 배경 + 파란 하단 라인)"""
    rect(slide, 0, 0, 13.33, 0.73, fill=NAVY)
    # 파란 하단 강조선
    rect(slide, 0, 0.73, 7.0, 0.03, fill=BLUE)
    rect(slide, 7.0, 0.73, 6.33, 0.03,
         fill=rgb(0xD0, 0xD8, 0xE8))
    # 좌측 흰 세로 막대
    rect(slide, 0.5, 0.13, 0.04, 0.48, fill=rgb(0xFF,0xFF,0xFF))
    txt(slide, 0.6, 0.14, 10.5, 0.45, title, 18, bold=True, color=WHITE)
    txt(slide, 11.8, 0.22, 1.2, 0.3,
        f"{slide_num:02d} / 08", 10, color=rgb(0xAA,0xBB,0xCC),
        align=PP_ALIGN.RIGHT)

def card_rect(slide, x, y, w, h, fill=CARD, border=BORD, bw=0.75):
    """카드 도형"""
    return rect(slide, x, y, w, h, fill=fill, line=border, line_w=bw)

def colored_bar(slide, x, y, w, h, color):
    """컬러 상단/측면 막대"""
    return rect(slide, x, y, w, h, fill=color)

def lighten(c, amt=210):
    """RGBColor를 밝게 (배경색용) — hex str 파싱"""
    h = str(c)  # e.g. "1868C8"
    r = min(int(h[0:2], 16) + amt, 255)
    g = min(int(h[2:4], 16) + amt, 255)
    b = min(int(h[4:6], 16) + amt, 255)
    return RGBColor(r, g, b)


# ════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ════════════════════════════════════════════════════
s1 = add_slide()

# 상단 navy 절반
rect(s1, 0, 0, 13.33, 4.0, fill=NAVY)
# 파란 강조선
rect(s1, 0.5, 1.1, 9.0, 0.02, fill=rgb(0x40,0x80,0xC8))
# 하단 라이트 배경
rect(s1, 0, 4.0, 13.33, 3.5, fill=LIGHT)

# 로고 박스
card_rect(s1, 0.5, 0.28, 2.4, 0.55, fill=rgb(0x1E,0x40,0x70),
          border=rgb(0x40,0x70,0xA8), bw=1)
txt(s1, 0.55, 0.31, 2.3, 0.23, "대한기계설비산업연구원", 9,
    bold=True, color=rgb(0xC0,0xD8,0xFF))
txt(s1, 0.55, 0.52, 2.3, 0.22, "KRIMFI", 8.5,
    color=rgb(0x70,0x90,0xB8))

# 컨퍼런스 배지
card_rect(s1, 10.5, 0.32, 2.4, 0.42, fill=rgb(0x1E,0x40,0x70),
          border=rgb(0x40,0x70,0xA8), bw=1)
txt(s1, 10.5, 0.42, 2.4, 0.25, "2026 학술발표대회", 9.5,
    color=rgb(0xC0,0xD8,0xFF), align=PP_ALIGN.CENTER)

# 제목
txt(s1, 1.0, 1.2, 11.3, 1.5,
    "기계설비 성능점검", 42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s1, 1.0, 2.55, 11.3, 0.8,
    "검토자문 시스템 개발", 42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 영문 부제
txt(s1, 1.0, 3.28, 11.3, 0.4,
    "Development of a System for Reviewing and Advising on Mechanical Equipment Performance Inspections",
    12, color=rgb(0x90,0xB0,0xD8), align=PP_ALIGN.CENTER)

# 구분선
rect(s1, 4.5, 3.82, 4.33, 0.01, fill=rgb(0x40,0x80,0xC8))

# 저자 카드
card_rect(s1, 3.2, 4.12, 6.93, 0.95, fill=CARD, border=BORD, bw=1)
txt(s1, 3.2, 4.22, 6.93, 0.35,
    "이창재¹   양자강¹   이영준²†", 14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
txt(s1, 3.2, 4.55, 6.93, 0.3,
    "¹대한기계설비산업연구원 선임연구원   ²대한기계설비산업연구원 연구위원",
    10.5, color=TEXT2, align=PP_ALIGN.CENTER)

# 하단 키워드
txt(s1, 1.0, 7.12, 11.3, 0.28,
    "Key words: 기계설비 성능점검  ·  기계설비법  ·  검토자문시스템",
    10, color=TEXT3, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 2 — 목차
# ════════════════════════════════════════════════════
s2 = add_slide()
rect(s2, 0, 0, 13.33, 7.5, fill=LIGHT)
hdr(s2, "목  차", 2)

TOC = [
    ("01", "연구 배경 및 필요성",    "기계설비 성능점검 제도 현황 및 문제점",    True),
    ("02", "검토자문 프로세스 설계", "4개 주체 역할분담 및 다단계 검토 흐름",   False),
    ("03", "시스템 구현",            "웹 기반 플랫폼 주요 화면",                False),
    ("04", "검토항목 구성",          "보고서 적정성 평가 4대 항목",              False),
    ("05", "기대효과 및 적용 현황",  "도입 효과 및 서울시 공공건축물 시범 적용", False),
    ("06", "결론 및 향후 연구",      "연구 성과 요약 및 후속 연구 계획",         False),
]

cols = [(0.45, 6.22), (6.67, 6.22)]
rows = 3
item_h = (7.5 - 0.90) / rows - 0.12

for i, (num, title, sub, active) in enumerate(TOC):
    col = i % 2
    row = i // 2
    x = cols[col][0]
    y = 0.86 + row * (item_h + 0.12)
    w = cols[col][1]

    if active:
        # 활성 항목: 파란 번호 박스
        card_rect(s2, x, y, w, item_h, fill=CARD, border=BORD, bw=1)
        rect(s2, x, y, 0.65, item_h, fill=BLUE)
        txt(s2, x, y + item_h/2 - 0.2, 0.65, 0.4, num, 20,
            bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    else:
        card_rect(s2, x, y, w, item_h, fill=CARD, border=BORD, bw=1)
        rect(s2, x, y, 0.65, item_h, fill=BG3)
        # 좌측 파란 테두리
        rect(s2, x + 0.61, y, 0.04, item_h, fill=BLUE)
        txt(s2, x + 0.03, y + item_h/2 - 0.2, 0.6, 0.4, num, 20,
            bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    txt(s2, x + 0.73, y + 0.10, w - 0.8, 0.35, title, 13, bold=True, color=TEXT)
    txt(s2, x + 0.73, y + 0.42, w - 0.8, 0.28, sub, 10.5, color=TEXT2)


# ════════════════════════════════════════════════════
# SLIDE 3 — 연구 배경 및 필요성
# ════════════════════════════════════════════════════
s3 = add_slide()
rect(s3, 0, 0, 13.33, 7.5, fill=LIGHT)
hdr(s3, "연구 배경 및 필요성", 3)

# 법령 배너
rect(s3, 0.45, 0.86, 12.43, 0.65, fill=NAVY)
txt(s3, 0.7, 0.91, 11.5, 0.25,
    "「기계설비법」 제17조  →  일정 규모 이상 건축물 관리주체: 매년 1회 이상 성능점검 의무",
    12, bold=False, color=rgb(0xC8,0xDC,0xF4))
txt(s3, 0.7, 1.12, 11.5, 0.25,
    "2021년 8월 시행  ·  연면적 규모별 순차 확대 적용 중  ·  2026년 현재 약 5년 시행",
    10.5, color=rgb(0x80,0xA0,0xC8))

# 4 문제점 카드
PROBS = [
    (ORANGE, "①", "보고서 품질 편차",
     "업체 간 기술역량 차이\n점검항목 해석 상이\n보고서 서식 비표준화\n내용 구성 불균일"),
    (BLUE, "②", "검토 역량 한계",
     "지자체 담당 공무원\n기술 전문성 부족\n행정적 여력 한계\n개별 검토 불가"),
    (PURPLE, "③", "검토 수요 급증",
     "성능점검 대상 건축물\n지속 확대 시행\n연간 제출 보고서 수\n증가 추세"),
    (GREEN, "④", "전문 검토 필요",
     "현장 실무자 및\n전문가 의견에서\n전문적 검토 과정의\n필요성 지속 제기"),
]
card_w = (12.43 - 3*0.15) / 4
for i, (col, num, title, body) in enumerate(PROBS):
    x = 0.45 + i * (card_w + 0.15)
    card_rect(s3, x, 1.6, card_w, 4.2, fill=CARD, border=BORD, bw=1)
    rect(s3, x, 1.6, card_w, 0.9, fill=col)
    txt(s3, x + 0.1, 1.63, card_w - 0.2, 0.35, num, 22,
        bold=True, color=rgb(0xFF,0xFF,0xFF))
    txt(s3, x + 0.1, 1.92, card_w - 0.2, 0.4, title, 13,
        bold=True, color=WHITE)
    for j, line in enumerate(body.split('\n')):
        txt(s3, x + 0.15, 2.6 + j*0.35, card_w - 0.25, 0.32,
            line, 11.5, color=TEXT2)

# 결론 배너
rect(s3, 0.45, 5.95, 12.43, 0.56, fill=BLUEL,
     line=BLUE, line_w=1.2)
txt(s3, 0.6, 6.02, 0.25, 0.38, "→", 18, bold=True, color=BLUE)
txt(s3, 0.9, 6.03, 11.5, 0.38,
    "성능점검 결과보고서의 품질 제고 및 지자체 담당 공무원의 보고서 적정성 판단 지원을 위한 검토자문 시스템 개발",
    12.5, bold=True, color=NAVY)


# ════════════════════════════════════════════════════
# SLIDE 4 — 검토자문 프로세스 설계
# ════════════════════════════════════════════════════
s4 = add_slide()
rect(s4, 0, 0, 13.33, 7.5, fill=LIGHT)
hdr(s4, "검토자문 프로세스 설계", 4)

# 설명 텍스트
txt(s4, 0.45, 0.84, 12.0, 0.3,
    "4개 주체 간 역할 분담 기반 다단계 검토 구조 — 단일 검토자 판단 편향 방지 + 보고서 품질 체계적 확보",
    11.5, color=TEXT2)

# 컬럼 헤더 (4개)
COLS4 = [
    (0.45, 2.78, NAVY,   "관리주체\nBuilding Owner"),
    (3.38, 2.78, BLUE,   "성능점검업자\nInspection Contractor"),
    (6.31, 2.78, NAVY2,  "연구원 (KRIMFI)"),
    (9.24, 3.64, PURPLE, "검토자문위원단\nReview Advisory Committee"),
]
for x, w, col, label in COLS4:
    rect(s4, x, 1.12, w, 0.52, fill=col)
    lines = label.split('\n')
    txt(s4, x, 1.15, w, 0.28, lines[0], 13,
        bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if len(lines) > 1:
        txt(s4, x, 1.38, w, 0.22, lines[1], 9,
            color=rgb(0xCC,0xDD,0xEE), align=PP_ALIGN.CENTER)
    # 컬럼 배경
    rect(s4, x, 1.64, w, 5.6, fill=rgb(0xF8,0xFA,0xFC),
         line=rgb(0xE0,0xE8,0xF0), line_w=0.5)

# 박스 정의: (col_x, col_w, y, height, label_ko, label_en, fill, border)
BOX_X = [0.45, 3.38, 6.31, 9.24]
BOX_W = [2.78, 2.78, 2.78, 3.64]
BOX_H = 0.52

def flow_box(slide, col, y, label_ko, label_en=None,
             fill=CARD, border=BORD, bw=0.75):
    x = BOX_X[col]
    w = BOX_W[col]
    card_rect(slide, x + 0.12, y, w - 0.24, BOX_H,
              fill=fill, border=border, bw=bw)
    txt(slide, x + 0.12, y + 0.04, w - 0.24, 0.28,
        label_ko, 12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    if label_en:
        txt(slide, x + 0.12, y + 0.3, w - 0.24, 0.2,
            label_en, 9, color=TEXT3, align=PP_ALIGN.CENTER)

Y1, Y2, Y3, Y4, Y5 = 1.72, 2.42, 3.12, 3.82, 4.65

flow_box(s4, 0, Y1, "성능점검 발주", "Outsource Performance Inspection")
flow_box(s4, 1, Y1, "입찰·낙찰·계약", "Bidding & Contract")
flow_box(s4, 1, Y2, "성능점검 수행", "Perform Performance Inspection")
flow_box(s4, 1, Y3, "결과보고서 작성", "Inspection Report Preparation")
flow_box(s4, 2, Y3, "보고서 접수", "Report Reception",
         fill=BLUEL, border=BLUE, bw=1.2)
flow_box(s4, 3, Y3, "1차 검토·의견 작성", "1st Review & Write Opinion",
         fill=rgb(0xF0,0xED,0xF8), border=PURPLE, bw=1.2)
flow_box(s4, 2, Y4, "2차 검토", "2nd Review (Opinion Check)",
         fill=BLUEL, border=BLUE, bw=1.2)
flow_box(s4, 1, Y4, "보완 회신 (환류)", "Revision Response",
         fill=rgb(0xFD,0xF0,0xEB), border=ORANGE, bw=1)
flow_box(s4, 3, Y4, "최종 검토·적합 판정", "Final Review & Approved",
         fill=rgb(0xF0,0xED,0xF8), border=PURPLE, bw=1.2)
flow_box(s4, 2, Y5, "검토확인서 발급", "Review Certificate Issued",
         fill=GREENL, border=GREEN, bw=1.5)
flow_box(s4, 1, Y5, "검토확인서 수령", "Certificate Receipt",
         fill=GREENL, border=GREEN, bw=1.2)
flow_box(s4, 0, Y5, "보고서·확인서 수령", "Receive Report & Certificate",
         fill=GREENL, border=GREEN, bw=1.2)

# 지자체 제출
rect(s4, BOX_X[0] + 0.12, 5.3, BOX_W[0] - 0.24, 0.52, fill=NAVY)
txt(s4, BOX_X[0] + 0.12, 5.33, BOX_W[0] - 0.24, 0.28,
    "지자체(MIS) 제출", 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s4, BOX_X[0] + 0.12, 5.59, BOX_W[0] - 0.24, 0.2,
    "Submit to Local Gov't (MIS)", 9, color=rgb(0xAA,0xCC,0xEE),
    align=PP_ALIGN.CENTER)

# 화살표 (단순 텍스트 ↓ / → 사용)
ARROWS = [
    # (x, y, w, h, text, size, color, bold)
    (1.84, 2.20, 0.3, 0.22, "↓", 14, TEXT3, False),
    (1.84, 2.90, 0.3, 0.22, "↓", 14, TEXT3, False),
    (6.16, 3.34, 0.2, 0.22, "→", 12, BLUE,  True),
    (9.09, 3.34, 0.2, 0.22, "→", 12, PURPLE, True),
    (6.16, 4.04, 0.2, 0.22, "→", 12, PURPLE, True),
    (6.13, 4.54, 0.15, 0.22, "→", 12, ORANGE, True),
    (9.09, 4.04, 0.2, 0.22, "→", 12, PURPLE, True),
    (6.13, 4.88, 0.2, 0.22, "→", 12, GREEN, True),
    (3.25, 4.88, 0.15, 0.22, "→", 12, GREEN, True),
    (1.84, 5.17, 0.3, 0.22, "↓", 14, NAVY,  True),
]
for (ax, ay, aw, ah, at, asz, acol, abold) in ARROWS:
    txt(s4, ax, ay, aw, ah, at, asz, bold=abold, color=acol, align=PP_ALIGN.CENTER)

# 범례
LEG = [
    (0.45, GREEN,  "합격·발급 흐름"),
    (2.3,  BLUE,   "접수·검토"),
    (4.15, PURPLE, "검토위원회"),
    (6.0,  ORANGE, "보완 요청·환류"),
]
for lx, lc, lt in LEG:
    rect(s4, lx, 6.9, 0.25, 0.14, fill=lc)
    txt(s4, lx + 0.3, 6.88, 1.6, 0.2, lt, 10, color=TEXT2)


# ════════════════════════════════════════════════════
# SLIDE 5 — 시스템 구현
# ════════════════════════════════════════════════════
s5 = add_slide()
rect(s5, 0, 0, 13.33, 7.5, fill=LIGHT)
hdr(s5, "시스템 구현", 5)

# 태그 텍스트
TAGS = [
    (BLUE,   "웹 기반 온라인 플랫폼"),
    (GREEN,  "전 과정 통합 관리"),
    (NAVY,   "실시간 진행 상태 확인"),
    (PURPLE, "이력 관리·투명성 확보"),
]
tx = 0.45
for tc, tl in TAGS:
    rect(s5, tx, 0.84, 1.95, 0.3, fill=lighten(tc, 220),
         line=tc, line_w=0.75)
    txt(s5, tx + 0.06, 0.86, 1.85, 0.26, tl, 10.5,
        bold=True, color=tc, align=PP_ALIGN.CENTER)
    tx += 2.05

# 3 화면 캡처 플레이스홀더
SCREENS = [
    (NAVY,  "Fig. 2", "보고서 접수현황 화면",      "Report management screen",
     "[ 화면 캡처 삽입 ]\n(보고서 목록·검토 진행 단계)",
     "건축물 정보 · 점검업체 · 접수일자\n검토 단계: 접수 / 검토 중 / 보완요청 / 완료"),
    (PURPLE,"Fig. 3", "검토의견서 작성 화면",       "Review opinion form screen",
     "[ 화면 캡처 삽입 ]\n(검토자문위원 작성 화면)",
     "항목별 적정성 평가\n구체적 검토의견 기술"),
    (GREEN, "Fig. 4", "의견작성 확인 화면 (연구원)", "Report opinion confirmation screen",
     "[ 화면 캡처 삽입 ]\n(연구원 2차 확인 화면)",
     "1차 검토의견 일관성·객관성 확인\n보완·조정 후 승인요청"),
]
sw = (12.43 - 2*0.15) / 3
for i, (col, fig, ko, en, ph, desc) in enumerate(SCREENS):
    x = 0.45 + i * (sw + 0.15)
    # 헤더
    rect(s5, x, 1.24, sw, 0.65, fill=col)
    txt(s5, x + 0.1, 1.27, sw - 0.2, 0.2, fig, 9,
        color=rgb(0xFF,0xFF,0xFF), bold=False)
    txt(s5, x + 0.1, 1.45, sw - 0.2, 0.28, ko, 12.5,
        bold=True, color=WHITE)
    # 플레이스홀더
    card_rect(s5, x, 1.89, sw, 3.0, fill=rgb(0xF8,0xFB,0xFF),
              border=rgb(0xB0,0xC8,0xE8), bw=1)
    txt(s5, x + 0.1, 2.9, sw - 0.2, 0.7, ph, 11,
        color=TEXT3, align=PP_ALIGN.CENTER)
    # 설명
    card_rect(s5, x, 4.89, sw, 0.7, fill=BG3, border=BORD, bw=0.5)
    for j, dl in enumerate(desc.split('\n')):
        txt(s5, x + 0.12, 4.94 + j*0.3, sw - 0.24, 0.28, dl, 11,
            color=TEXT2)

# 플랫폼 특징 박스
rect(s5, 0.45, 5.72, 12.43, 0.28, fill=NAVY)
txt(s5, 0.55, 5.75, 5.0, 0.25, "플랫폼 특징", 12,
    bold=True, color=WHITE)
FEATS = [
    "성능점검업자 온라인 보고서 접수",
    "검토 진행 단계 실시간 확인",
    "보완 요청·재검토 환류 구조",
    "검토자문위원 1차 검토의견 작성",
    "연구원 2차 검토 · 최종 승인요청",
    "검토확인서 자동 작성·발급",
]
card_rect(s5, 0.45, 6.0, 12.43, 1.3, fill=CARD, border=BORD, bw=0.75)
for i, f in enumerate(FEATS):
    col_x = 0.6 + (i % 3) * 4.1
    row_y = 6.08 + (i // 3) * 0.42
    txt(s5, col_x, row_y, 3.9, 0.35, f"▸  {f}", 11.5, color=TEXT2)


# ════════════════════════════════════════════════════
# SLIDE 6 — 검토항목 구성
# ════════════════════════════════════════════════════
s6 = add_slide()
rect(s6, 0, 0, 13.33, 7.5, fill=LIGHT)
hdr(s6, "검토항목 구성", 6)

txt(s6, 0.45, 0.84, 12.0, 0.3,
    "검토자문위원은 성능점검 결과보고서의 항목별 적정성을 평가하고 구체적인 검토의견 기술 — 4개 주요 검토 카테고리",
    11.5, color=TEXT2)

ITEMS = [
    (BLUE,   "1", "점검대상 기계설비 현황의 정확성",
     ["설치된 기계설비의 종류·규모·수량의 정확한 기재 여부",
      "설비 제원 및 용량 정보의 설계도서 일치 여부",
      "점검 대상 범위의 적정성"]),
    (PURPLE, "2", "성능점검표 기재 내용의 충실도",
     ["점검항목 전체에 대한 빠짐없는 기재 여부",
      "점검 결과 수치의 구체적 기록 여부",
      "측정 조건 및 방법의 명시 여부"]),
    (ORANGE, "3", "설계값 대비 측정값의 적정 범위 여부",
     ["설계값 대비 측정값 비율의 허용 범위 내 여부",
      "측정 기준 단위 및 환산의 적정성",
      "기준 초과 항목의 원인 분석 기재 여부"]),
    (GREEN,  "4", "성능개선계획의 타당성",
     ["불합격 항목에 대한 개선계획 수립 여부",
      "개선 방법·일정·담당자의 명확한 기재",
      "개선 조치의 기술적 타당성"]),
]

iw = (12.43 - 0.15) / 2
ih = (7.5 - 0.86 - 0.7 - 0.15) / 2 - 0.1
positions = [
    (0.45, 1.22), (0.45 + iw + 0.15, 1.22),
    (0.45, 1.22 + ih + 0.12), (0.45 + iw + 0.15, 1.22 + ih + 0.12),
]

for i, ((ix, iy), (col, num, title, bullets)) in enumerate(zip(positions, ITEMS)):
    card_rect(s6, ix, iy, iw, ih, fill=CARD, border=BORD, bw=1)
    rect(s6, ix, iy, 0.1, ih, fill=col)
    # 번호 원
    rect(s6, ix + 0.2, iy + 0.15, 0.42, 0.42, fill=lighten(col, 200),
         line=col, line_w=1.5)
    txt(s6, ix + 0.2, iy + 0.2, 0.42, 0.3, num, 16,
        bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s6, ix + 0.72, iy + 0.18, iw - 0.85, 0.35,
        title, 13.5, bold=True, color=TEXT)
    for j, b in enumerate(bullets):
        txt(s6, ix + 0.25, iy + 0.73 + j*0.42, iw - 0.38, 0.35,
            f"●  {b}", 11.5, color=TEXT2)

# 결과 판정 배너
rect(s6, 0.45, 6.72, 3.9, 0.52, fill=GREENL, line=GREEN, line_w=1)
txt(s6, 0.5, 6.76, 0.3, 0.35, "✓", 18, bold=True, color=GREEN)
txt(s6, 0.85, 6.8, 3.0, 0.28, "합격 → 검토확인서 발급 → 지자체 제출",
    11, color=NAVY)

rect(s6, 4.5, 6.72, 4.3, 0.52,
     fill=rgb(0xFD,0xF0,0xEB), line=ORANGE, line_w=1)
txt(s6, 4.56, 6.76, 0.3, 0.35, "⚠", 18, bold=True, color=ORANGE)
txt(s6, 4.9, 6.8, 3.8, 0.28,
    "보완 요청 → 성능점검업자 보완 회신 → 재검토",
    11, color=NAVY)

rect(s6, 8.95, 6.72, 3.93, 0.52,
     fill=rgb(0xFD,0xF0,0xEB), line=rgb(0xC0,0x50,0x10), line_w=1)
txt(s6, 9.0, 6.76, 0.3, 0.35, "✗", 18, bold=True, color=ORANGE)
txt(s6, 9.35, 6.8, 3.4, 0.28,
    "반려 → 반려 사유 통보",
    11, color=NAVY)


# ════════════════════════════════════════════════════
# SLIDE 7 — 기대효과 및 적용 현황
# ════════════════════════════════════════════════════
s7 = add_slide()
rect(s7, 0, 0, 13.33, 7.5, fill=LIGHT)
hdr(s7, "기대효과 및 적용 현황", 7)

EFFECTS = [
    (NAVY,  "효과 1", "보고서 품질 균질화",
     "전문 검토자문위원단에 의한\n다단계 검토 체계 적용\n\n→ 성능점검 결과보고서의\n   품질 수준 균질화"),
    (BLUE,  "효과 2", "행정 효율성 향상",
     "지자체 공무원은\n검토확인서 기반으로\n보고서 적정 여부 신속 판단\n\n→ 행정 처리 시간 단축"),
    (GREEN, "효과 3", "업체 역량 점진적 개선",
     "보완 요청 및 재검토\n환류 구조를 통해\n성능점검업체의 보고서\n작성 역량 점진적 개선"),
]
ew = (12.43 - 2*0.15) / 3
for i, (col, badge, title, body) in enumerate(EFFECTS):
    x = 0.45 + i * (ew + 0.15)
    card_rect(s7, x, 0.86, ew, 2.9, fill=CARD, border=BORD, bw=1)
    rect(s7, x, 0.86, ew, 0.7, fill=col)
    txt(s7, x + 0.1, 0.88, 1.2, 0.2, badge, 9.5,
        color=rgb(0xFF,0xFF,0xFF))
    txt(s7, x + 0.1, 1.06, ew - 0.2, 0.35, title, 15,
        bold=True, color=WHITE)
    for j, bl in enumerate(body.split('\n')):
        txt(s7, x + 0.15, 1.72 + j * 0.35, ew - 0.25, 0.32,
            bl, 11.5, color=TEXT2)

# 적용 현황 카드
rect(s7, 0.45, 3.88, 12.43, 0.34, fill=BG3, line=BLUE, line_w=1.5)
txt(s7, 0.6, 3.92, 4.0, 0.28, "시스템 적용 현황", 13,
    bold=True, color=NAVY)

# 타임라인 박스들
TL = [
    (NAVY2,  "2021.08", "기계설비법 시행"),
    (BLUE,   "2025~",   "시스템 개발"),
    (GREEN,  "2026.04~","서울시 공공건축물 적용"),
    (PURPLE, "계획",    "민간 건축물 확대"),
    (NAVY,   "후속",    "AI 사전검토 도입"),
]
tw = (12.43 - 4*0.2) / 5
for i, (col, date, desc) in enumerate(TL):
    x = 0.45 + i * (tw + 0.2)
    rect(s7, x + tw/2 - 0.3, 4.35, 0.6, 0.6, fill=lighten(col, 200),
         line=col, line_w=1.5)
    txt(s7, x, 4.38, tw, 0.3, ["⚖","🔧","✅","🏢","🤖"][i], 18,
        align=PP_ALIGN.CENTER)
    txt(s7, x, 5.1, tw, 0.28, date, 10, bold=True,
        color=col, align=PP_ALIGN.CENTER)
    txt(s7, x, 5.38, tw, 0.28, desc, 9.5,
        color=TEXT2, align=PP_ALIGN.CENTER)
    # 연결선 (마지막 제외)
    if i < 4:
        rect(s7, x + tw - 0.01, 4.62, 0.22, 0.01,
             fill=rgb(0xB0,0xC0,0xD8))

# 현황 하이라이트
rect(s7, 0.45, 5.88, 12.43, 0.56, fill=GREENL, line=GREEN, line_w=1.2)
txt(s7, 0.6, 5.94, 0.28, 0.38, "📍", 16)
txt(s7, 0.96, 5.97, 11.5, 0.35,
    "현재 적용 중: 2026년 4월부터 서울특별시 공공건축물을 대상으로 우선 적용  —  향후 민간 건축물로 확대 예정",
    12.5, color=NAVY)


# ════════════════════════════════════════════════════
# SLIDE 8 — 결론 및 향후 연구
# ════════════════════════════════════════════════════
s8 = add_slide()
rect(s8, 0, 0, 13.33, 7.5, fill=LIGHT)
hdr(s8, "결론 및 향후 연구", 8)

# 결론 요약 카드
card_rect(s8, 0.45, 0.84, 12.43, 1.45, fill=CARD, border=BORD, bw=1)
rect(s8, 0.45, 0.84, 0.06, 1.45, fill=NAVY)
txt(s8, 0.6, 0.88, 2.5, 0.3, "연구 결론", 12.5,
    bold=True, color=NAVY)
txt(s8, 0.6, 1.17, 12.0, 1.0,
    "성능점검 결과보고서의 품질 편차 문제와 지방자치단체 담당 공무원의 검토 역량 한계를 해소하기 위해\n"
    "'기계설비 성능점검 검토자문 시스템'을 개발하였다. 보고서 온라인 접수, 검토자문위원단의 1차 검토,\n"
    "연구원의 2차 검토, 최종 판정 및 검토확인서 발급에 이르는 전 과정을 웹 기반 플랫폼으로 통합 구현하여\n"
    "검토자문 절차의 체계화와 이력관리의 투명성을 확보하였다.",
    11.5, color=TEXT2)

# 3 성과 카드
ACHS = [
    (NAVY,   "✓ 절차 체계화",
     "접수 → 1차 검토 → 2차 검토\n→ 최종 판정 → 확인서 발급\n전 과정 단일 플랫폼 통합"),
    (BLUE,   "✓ 전문성 확보",
     "기술사·박사급 전문 검토자문\n위원단 참여로 보고서 적정성\n전문적 판단 가능"),
    (GREEN,  "✓ 행정 지원",
     "검토확인서 기반 지자체\n행정 판단 지원\n이력관리 투명성 확보"),
]
aw = (12.43 - 2*0.15) / 3
for i, (col, title, body) in enumerate(ACHS):
    x = 0.45 + i * (aw + 0.15)
    card_rect(s8, x, 2.42, aw, 2.1, fill=CARD, border=BORD, bw=1)
    rect(s8, x, 2.42, aw, 0.04, fill=col)
    txt(s8, x + 0.15, 2.5, aw - 0.25, 0.35, title, 13,
        bold=True, color=col)
    for j, bl in enumerate(body.split('\n')):
        txt(s8, x + 0.15, 2.96 + j * 0.38, aw - 0.25, 0.32,
            bl, 11.5, color=TEXT2)

# 향후 연구
rect(s8, 0.45, 4.66, 12.43, 0.32, fill=BG3, line=BLUE, line_w=1.5)
txt(s8, 0.6, 4.7, 4.0, 0.26, "향후 연구 계획", 12.5,
    bold=True, color=NAVY)

FUTURES = [
    (BLUE,   "🤖", "AI 기반 사전검토",
     "축적된 검토 데이터를 활용한\nAI 기반 사전검토 기능 도입"),
    (GREEN,  "🔗", "MIS 연계 강화",
     "기계설비산업민원웹포털(MIS)\n연계 강화"),
    (PURPLE, "🗺", "타 지자체 확산",
     "서울 이외 지방자치단체로의\n확산 방안 검토"),
]
fw = (12.43 - 2*0.15) / 3
for i, (col, icon, title, desc) in enumerate(FUTURES):
    x = 0.45 + i * (fw + 0.15)
    card_rect(s8, x, 5.06, fw, 1.22, fill=CARD, border=BORD, bw=0.75)
    txt(s8, x + 0.12, 5.1, 0.4, 0.35, icon, 18)
    txt(s8, x + 0.55, 5.12, fw - 0.65, 0.32, title, 12.5,
        bold=True, color=col)
    for j, dl in enumerate(desc.split('\n')):
        txt(s8, x + 0.12, 5.52 + j * 0.3, fw - 0.2, 0.28,
            dl, 10.5, color=TEXT2)

# 감사
rect(s8, 0.45, 6.4, 12.43, 0.88, fill=NAVY)
txt(s8, 0, 6.52, 13.33, 0.38, "감사합니다", 22,
    bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s8, 0, 6.88, 13.33, 0.3,
    "본 연구는 대한기계설비산업연구원의 지원으로 수행되었습니다.   |   leeyj@krimfi.re.kr",
    10, color=rgb(0xAA,0xCC,0xEE), align=PP_ALIGN.CENTER)


# ── 저장 ──────────────────────────────────────────
output = "검토자문시스템_발표자료.pptx"
prs.save(output)
print(f"[OK] 저장 완료: {output}")
