#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""기계설비 성능점검 검토자문 시스템 → PPTX 변환"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── 단위 ──────────────────────────────────────────────
PX = 9525          # 1px = 9525 EMU  (96dpi 기준)

def px(n): return int(n * PX)

# ── 색상 ──────────────────────────────────────────────
NAVY  = RGBColor(0x15, 0x40, 0x63)
GREEN = RGBColor(0x73, 0xBF, 0x45)
RED   = RGBColor(0xEF, 0x41, 0x29)
PAPER = RGBColor(0xf4, 0xf6, 0xf9)
INK   = RGBColor(0x10, 0x24, 0x3a)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6a, 0x7a, 0x8c)
LBLUE = RGBColor(0xc6, 0xd6, 0xe6)
DFOOT = RGBColor(0x9d, 0xb6, 0xcf)
PLINE = RGBColor(0xdd, 0xe3, 0xea)
GFILL = RGBColor(0x3f, 0x7a, 0x1e)
RFILL = RGBColor(0xb3, 0x2a, 0x18)

FONT = "맑은 고딕"

# ── 프레젠테이션 ──────────────────────────────────────
prs = Presentation()
prs.slide_width  = px(1280)
prs.slide_height = px(720)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]   # 빈 레이아웃

# ── 공통 유틸 ─────────────────────────────────────────

def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def tb(slide, text, l, t, w, h,
       size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    """텍스트박스 하나 추가"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text           = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.color.rgb = color
    r.font.name      = FONT
    return box


def tb2(slide, run_specs, l, t, w, h, align=PP_ALIGN.LEFT, wrap=True):
    """
    여러 스타일의 run을 가진 텍스트박스.
    run_specs: [(text, size, bold, color), ...]
    """
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    for text, size, bold, color in run_specs:
        r = p.add_run()
        r.text           = text
        r.font.size      = Pt(size)
        r.font.bold      = bold
        r.font.color.rgb = color
        r.font.name      = FONT
    return box


def no_line(shape):
    """도형 테두리 제거"""
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    try:
        sp_pr = shape._element.spPr
    except AttributeError:
        NS_P  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        sp_pr = shape._element.find(f'{{{NS_P}}}spPr')
    if sp_pr is None:
        return
    NS_A_URI = f'{{{NS_A}}}'
    ln = sp_pr.find(f'{NS_A_URI}ln')
    if ln is not None:
        sp_pr.remove(ln)
    ln = etree.SubElement(sp_pr, f'{NS_A_URI}ln')
    etree.SubElement(ln, f'{NS_A_URI}noFill')


def add_rect(slide, l, t, w, h, fill, border=None, bw=Pt(1.5)):
    """직사각형 도형"""
    s = slide.shapes.add_shape(1, l, t, w, h)   # 1 = rect
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        no_line(s)
    return s


def add_rounded_rect(slide, l, t, w, h, fill, border=None, bw=Pt(1.5)):
    """모서리 둥근 직사각형"""
    s = slide.shapes.add_shape(5, l, t, w, h)   # 5 = roundRect
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        no_line(s)
    return s


# ── 슬라이드 공통 요소 ────────────────────────────────

def add_eyebrow(slide, text, dark=True):
    tb(slide, text, px(88), px(74), W - px(176), px(36),
       size=13, bold=True, color=GREEN if dark else NAVY)


def add_h2(slide, text, dark=False):
    tb(slide, text, px(88), px(118), W - px(176), px(108),
       size=38, bold=True, color=WHITE if dark else INK)


def add_footer(slide, left_t, right_t='', dark=True):
    c = DFOOT if dark else MUTED
    tb(slide, left_t,  px(88),        H - px(60), px(700), px(40), size=12, color=c)
    tb(slide, right_t, W - px(300),   H - px(60), px(212), px(40), size=12, color=c,
       align=PP_ALIGN.RIGHT)


def add_bullet_row(slide, dot_char, item_text, y):
    """녹색 점(혹은 번호) + 항목 텍스트 한 줄"""
    add_rounded_rect(slide, px(88), y + px(4), px(34), px(34), GREEN)
    tb(slide, dot_char,
       px(88), y + px(2), px(34), px(38),
       size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
    tb(slide, item_text,
       px(140), y, W - px(228), px(68),
       size=20, color=INK)


def add_bullets(slide, items, sy=None):
    if sy is None:
        sy = px(242)
    for k, (d, t) in enumerate(items):
        add_bullet_row(slide, d, t, sy + k * px(74))


# ═══════════════════════════════════════════════════════
# 슬라이드 1: 표지  (dark)
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)

add_eyebrow(s, "대한기계설비산업연구원", dark=True)

tb(s, "기계설비 성능점검\n검토자문 시스템",
   px(88), px(128), W - px(176), px(256),
   size=62, bold=True, color=WHITE)

add_rect(s, px(88), px(392), px(78), px(6), GREEN)

tb(s, "AI 기반 성능점검 보고서 적정성 자동 검토",
   px(88), px(414), px(780), px(48), size=20, color=LBLUE)

add_footer(s, "KRIMFI", "2026", dark=True)


# ═══════════════════════════════════════════════════════
# 슬라이드 2: 배경 및 필요성  (light)
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, PAPER)
add_eyebrow(s, "배경 및 필요성", dark=False)
add_h2(s, "성능점검 보고서, 적정성 검토가 관건")
add_bullets(s, [
    ("·", "점검업체가 점검기준에 맞게 점검·작성했는지 확인이 필요"),
    ("·", "항목이 많고 계측 사진·수치가 혼재해 수작업 검토 부담이 큼"),
    ("·", "검토자에 따라 판정이 달라 일관성 확보가 어려움"),
])


# ═══════════════════════════════════════════════════════
# 슬라이드 3: 목적  (light)
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, PAPER)
add_eyebrow(s, "목적", dark=False)
add_h2(s, "검토를 자동화하고 기준을 표준화")
add_bullets(s, [
    ("·", "점검결과의 적정성을 AI가 검토해 적정 / 보완필요로 판정"),
    ("·", "측정값·설계값·비율을 포함한 정량적·기술적 자문의견 제공"),
    ("·", '"양호·적합" 같은 정성 표현에 그치지 않는 근거 중심 검토'),
])


# ═══════════════════════════════════════════════════════
# 슬라이드 4: 시스템 프로세스  (light)
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, PAPER)
add_eyebrow(s, "시스템 프로세스", dark=False)
add_h2(s, "업로드에서 자문의견까지 4단계")
add_bullets(s, [
    ("1", "보고서 업로드 — PDF 또는 현장 사진(다중)"),
    ("2", "AI 멀티모달 분석 — 전 장비·점검항목 추출, 사진 수치 판독"),
    ("3", "검토원칙 적용 — 항목별 적정 / 보완필요 판정"),
    ("4", "결과 구조화 — 기본정보·항목별 의견, 화면 분류·엑셀 출력"),
])


# ═══════════════════════════════════════════════════════
# 슬라이드 5: 검토 원칙  (light)
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, PAPER)
add_eyebrow(s, "검토 원칙", dark=False)
add_h2(s, "정량 검토를 위한 5가지 원칙")
add_bullets(s, [
    ("1", "설계값 대비 측정값 비율 계산·명시"),
    ("2", "단위 환산 직접 수행 (예: 풍속 → 풍량 CMH)"),
    ("3", "운전조건 확인 (정격 / 일반운전 구분)"),
    ("4", "판정 근거를 정량 논리로 서술"),
    ("5", "보완 필요 시 구체적 개선방향 제시"),
], sy=px(228))


# ═══════════════════════════════════════════════════════
# 슬라이드 6: 결과·실증  (light)
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, PAPER)
add_eyebrow(s, "결과 · 실증", dark=False)
add_h2(s, "충무아트센터 — 123개 항목 검토")

# 왼쪽: 64건 (적정)
tb(s, "64",        px(88),  px(242), px(310), px(180), size=115, bold=True, color=NAVY)
tb(s, "건",        px(356), px(296), px(88),  px(120), size=46,  bold=True, color=GREEN)
tb(s, "적정 (52%)", px(88), px(454), px(420), px(44),  size=18,  color=MUTED)

# 오른쪽: 59건 (보완필요)
tb(s, "59",          px(640), px(242), px(310), px(180), size=115, bold=True, color=RED)
tb(s, "건",          px(908), px(296), px(88),  px(120), size=46,  bold=True, color=RED)
tb(s, "보완필요 (48%)", px(640), px(454), px(420), px(44), size=18, color=MUTED)

add_footer(s, "대상 장비: 공기조화기 AHU-5·10·15·17, HVU-1, 팬코일 등", "", dark=False)


# ═══════════════════════════════════════════════════════
# 슬라이드 7: 대표 검토 사례  (light)
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, PAPER)
add_eyebrow(s, "대표 검토 사례", dark=False)
add_h2(s, "AHU-5 송풍기 풍량 — '적합'을 '보완필요'로 정정")

CL = px(88)
CR = px(676)
CT = px(226)
CW = px(508)
CH = px(364)
PAD = px(24)

# 왼쪽 카드 — 업체 점검결과
add_rect(s, CL, CT, CW, CH, WHITE, border=PLINE, bw=Pt(1))
tb(s, "업체 점검결과",
   CL + PAD, CT + PAD, CW - PAD*2, px(30),
   size=11, bold=True, color=RFILL)
tb(s, "적합",
   CL + PAD, CT + px(58), CW - PAD*2, px(54),
   size=24, bold=True, color=INK)
tb(s, "풍량 측정값에 대한 설계값 대비 정량 비교 없이\n'적합'으로 판정.",
   CL + PAD, CT + px(122), CW - PAD*2, px(190),
   size=15, color=MUTED, wrap=True)

# 오른쪽 카드 — AI 검토의견 (초록 테두리)
add_rect(s, CR, CT, CW, CH, WHITE, border=GREEN, bw=Pt(2))
tb(s, "AI 검토의견",
   CR + PAD, CT + PAD, CW - PAD*2, px(30),
   size=11, bold=True, color=GFILL)
tb(s, "보완필요",
   CR + PAD, CT + px(58), CW - PAD*2, px(54),
   size=24, bold=True, color=INK)
tb(s, "측정 18,634 CMH = 설계 21,000 CMH의 88.73%로\n기준(±10%, 90–110%) 불충족.\n정량 재판단 필요.",
   CR + PAD, CT + px(122), CW - PAD*2, px(210),
   size=15, color=MUTED, wrap=True)


# ═══════════════════════════════════════════════════════
# 슬라이드 8: 기대효과  (dark)
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)
add_eyebrow(s, "기대효과", dark=True)
add_h2(s, "검토 일관성 · 업무 효율 · 데이터 축적", dark=True)
tb(s, ("표준화된 기준으로 검토 품질을 균질화하고, 검토 업무를 효율화하며,\n"
       "건물 성능 데이터를 디지털로 축적·활용할 수 있습니다."),
   px(88), px(294), px(860), px(130), size=20, color=LBLUE, wrap=True)
add_footer(s, "감사합니다", "", dark=True)


# ── 저장 ──────────────────────────────────────────────
OUT = r"c:\Users\이창재\Desktop\py코딩연습\mpis-reviewer\기계설비_성능점검_검토자문시스템.pptx"
prs.save(OUT)
print(f"저장 완료: {OUT}")
