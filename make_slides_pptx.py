"""
기계설비 성능점검 검토자문 시스템 — 발표용 PPTX 생성 스크립트
실행: python make_slides_pptx.py
출력: 기계설비-성능점검-검토자문시스템.pptx (현재 디렉터리)
"""

import io
import os
import sys
import zipfile
import uuid

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# ── 색상 ──────────────────────────────────────────────────────
NAVY    = RGBColor(0x15, 0x40, 0x63)
NAVY_D  = RGBColor(0x0d, 0x2c, 0x46)
GREEN   = RGBColor(0x73, 0xBF, 0x45)
RED     = RGBColor(0xEF, 0x41, 0x29)
PAPER   = RGBColor(0xF4, 0xF6, 0xF9)
MUTED   = RGBColor(0x6A, 0x7A, 0x8C)
INK     = RGBColor(0x10, 0x24, 0x3A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_B = RGBColor(0xC6, 0xD6, 0xE6)

FONT    = "Pretendard"
W = Inches(13.33)
H = Inches(7.5)

# ── 폰트 경로 (pptx_maker.py 와 공유) ────────────────────────
FONT_DIR = os.path.join(os.path.dirname(__file__), "fonts", "pretendard")


# ── 헬퍼 ─────────────────────────────────────────────────────

def _set_para(para, align=PP_ALIGN.LEFT):
    para.alignment = align

def _run(para, text, size, bold=False, color=WHITE, italic=False):
    run = para.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run

def _textbox(slide, l, t, w, h, text, size, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    _set_para(p, align)
    _run(p, text, size, bold, color, italic)
    return txb

def _rect(slide, l, t, w, h, fill=None, line_color=None, line_w=0):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color and line_w:
        s.line.color.rgb = line_color; s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def _hdr(slide, title, subtitle=None):
    """공통 슬라이드 헤더 (네이비 바)"""
    _rect(slide, 0, 0, W, Inches(1.0), fill=NAVY)
    _textbox(slide, Inches(0.7), Inches(0.22), Inches(9), Inches(0.6),
             title, 22, bold=True, color=WHITE)
    if subtitle:
        _textbox(slide, Inches(0.7), Inches(0.75), Inches(11), Inches(0.4),
                 subtitle, 13, color=LIGHT_B)

def _footer(slide, text="KRIMFI", right="2026"):
    _textbox(slide, Inches(0.7), Inches(7.15), Inches(5), Inches(0.3),
             text, 11, color=MUTED)
    _textbox(slide, Inches(8.0), Inches(7.15), Inches(5), Inches(0.3),
             right, 11, color=MUTED, align=PP_ALIGN.RIGHT)

def _accent_rule(slide, l, t, w=Inches(1.0), h=Inches(0.07)):
    _rect(slide, l, t, w, h, fill=GREEN)

def _check_list(slide, items, start_y=Inches(1.4), gap=Inches(0.72),
                dot_color=GREEN, text_color=INK, size=20):
    for idx, text in enumerate(items):
        y = start_y + idx * gap
        # dot
        dot = slide.shapes.add_shape(1, Inches(0.7), y + Inches(0.04),
                                     Inches(0.42), Inches(0.42))
        dot.fill.solid(); dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        tf = dot.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = str(idx + 1) if text[0].isdigit() else '·'
        r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = RGBColor(0x0b, 0x2a, 0x12)
        # text
        _textbox(slide, Inches(1.3), y, Inches(11.5), Inches(0.5),
                 text.lstrip('1234567890').lstrip('·').strip(),
                 size, bold=False, color=text_color)


# ══════════════════════════════════════════════════════════════
# 슬라이드 1: 표지 (dark)
# ══════════════════════════════════════════════════════════════

def slide_cover(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    _rect(sl, 0, 0, W, H, fill=NAVY)
    # 상단 강조선
    _rect(sl, 0, 0, W, Inches(0.07), fill=GREEN)
    # 하단 강조선
    _rect(sl, 0, Inches(7.35), W, Inches(0.15), fill=GREEN)

    # 기관명
    _textbox(sl, Inches(1.5), Inches(1.4), Inches(10.3), Inches(0.55),
             "대한기계설비산업연구원", 15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # 메인 타이틀
    txb = sl.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.2))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for line in ["기계설비 성능점검", "검토자문 시스템"]:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.name = FONT; r.font.size = Pt(52); r.font.bold = True
        r.font.color.rgb = WHITE
        if line == "기계설비 성능점검":
            p = tf.add_paragraph()

    # 구분 룰
    _accent_rule(sl, Inches(5.5), Inches(4.35), Inches(2.33), Inches(0.06))

    # 부제
    _textbox(sl, Inches(1.5), Inches(4.55), Inches(10.3), Inches(0.6),
             "AI 기반 성능점검 보고서 적정성 자동 검토",
             18, color=LIGHT_B, align=PP_ALIGN.CENTER)

    _footer(sl, "KRIMFI", "2026")


# ══════════════════════════════════════════════════════════════
# 슬라이드 2: 배경 및 필요성 (light)
# ══════════════════════════════════════════════════════════════

def slide_background(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, W, H, fill=PAPER)
    _hdr(sl, "배경 및 필요성")
    _textbox(sl, Inches(0.7), Inches(1.05), Inches(12), Inches(0.7),
             "성능점검 보고서, 적정성 검토가 관건",
             32, bold=True, color=NAVY)
    items = [
        "점검업체가 점검기준에 맞게 점검·작성했는지 확인이 필요",
        "항목이 많고 계측 사진·수치가 혼재해 수작업 검토 부담이 큼",
        "검토자에 따라 판정이 달라 일관성 확보가 어려움",
    ]
    _check_list(sl, items, start_y=Inches(2.0), gap=Inches(0.95), text_color=INK)
    _footer(sl)


# ══════════════════════════════════════════════════════════════
# 슬라이드 3: 목적 (light)
# ══════════════════════════════════════════════════════════════

def slide_purpose(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, W, H, fill=PAPER)
    _hdr(sl, "목적")
    _textbox(sl, Inches(0.7), Inches(1.05), Inches(12), Inches(0.7),
             "검토를 자동화하고 기준을 표준화",
             32, bold=True, color=NAVY)
    items = [
        "점검결과의 적정성을 AI가 검토해 적정 / 보완필요로 판정",
        "측정값·설계값·비율을 포함한 정량적·기술적 자문의견 제공",
        "\"양호·적합\" 같은 정성 표현에 그치지 않는 근거 중심 검토",
    ]
    _check_list(sl, items, start_y=Inches(2.0), gap=Inches(0.95), text_color=INK)
    _footer(sl)


# ══════════════════════════════════════════════════════════════
# 슬라이드 4: 시스템 프로세스 (light)
# ══════════════════════════════════════════════════════════════

def slide_process(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, W, H, fill=PAPER)
    _hdr(sl, "시스템 프로세스")
    _textbox(sl, Inches(0.7), Inches(1.05), Inches(12), Inches(0.7),
             "업로드에서 자문의견까지 4단계",
             32, bold=True, color=NAVY)
    items = [
        "1  보고서 업로드 — PDF 또는 현장 사진(다중)",
        "2  AI 멀티모달 분석 — 전 장비·점검항목 추출, 사진 수치 판독",
        "3  검토원칙 적용 — 항목별 적정 / 보완필요 판정",
        "4  결과 구조화 — 기본정보·항목별 의견, 화면 분류·엑셀 출력",
    ]
    _check_list(sl, items, start_y=Inches(2.0), gap=Inches(0.88), text_color=INK)
    _footer(sl)


# ══════════════════════════════════════════════════════════════
# 슬라이드 5: 검토 원칙 (light)
# ══════════════════════════════════════════════════════════════

def slide_principles(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, W, H, fill=PAPER)
    _hdr(sl, "검토 원칙")
    _textbox(sl, Inches(0.7), Inches(1.05), Inches(12), Inches(0.7),
             "정량 검토를 위한 5가지 원칙",
             32, bold=True, color=NAVY)
    items = [
        "1  설계값 대비 측정값 비율 계산·명시",
        "2  단위 환산 직접 수행 (예: 풍속 → 풍량 CMH)",
        "3  운전조건 확인 (정격 / 일반운전 구분)",
        "4  판정 근거를 정량 논리로 서술",
        "5  보완 필요 시 구체적 개선방향 제시",
    ]
    _check_list(sl, items, start_y=Inches(1.95), gap=Inches(0.78), text_color=INK)
    _footer(sl)


# ══════════════════════════════════════════════════════════════
# 슬라이드 6: 결과·실증 (light)
# ══════════════════════════════════════════════════════════════

def slide_results(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, W, H, fill=PAPER)
    _hdr(sl, "결과 · 실증")
    _textbox(sl, Inches(0.7), Inches(1.05), Inches(12), Inches(0.7),
             "충무아트센터 — 123개 항목 검토",
             32, bold=True, color=NAVY)

    # 적정 블록
    _rect(sl, Inches(0.7), Inches(2.15), Inches(5.7), Inches(3.5), fill=PAPER)
    txb = sl.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(5.7), Inches(2.2))
    tf = txb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, "64", 110, bold=True, color=NAVY)
    _run(p, "건", 44, bold=True, color=GREEN)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    _run(p2, "적정 (52%)", 20, color=MUTED)

    # 보완필요 블록
    txb2 = sl.shapes.add_textbox(Inches(7.0), Inches(2.3), Inches(5.7), Inches(2.2))
    tf2 = txb2.text_frame
    p3 = tf2.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    _run(p3, "59", 110, bold=True, color=RED)
    _run(p3, "건", 44, bold=True, color=RED)
    p4 = tf2.add_paragraph(); p4.alignment = PP_ALIGN.CENTER
    _run(p4, "보완필요 (48%)", 20, color=MUTED)

    _textbox(sl, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
             "대상 장비: 공기조화기 AHU-5·10·15·17, HVU-1, 팬코일 등",
             12, color=MUTED)
    _footer(sl)


# ══════════════════════════════════════════════════════════════
# 슬라이드 7: 대표 검토 사례 (light)
# ══════════════════════════════════════════════════════════════

def slide_case(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, W, H, fill=PAPER)
    _hdr(sl, "대표 검토 사례")
    _textbox(sl, Inches(0.7), Inches(1.05), Inches(12), Inches(0.7),
             "AHU-5 송풍기 풍량 — '적합'을 '보완필요'로 정정",
             28, bold=True, color=NAVY)

    # 좌측 카드 (업체 - 빨강)
    lx, ly = Inches(0.7), Inches(2.0)
    cw, ch = Inches(5.8), Inches(4.3)
    _rect(sl, lx, ly, cw, ch, fill=WHITE, line_color=RED, line_w=1.5)
    _textbox(sl, lx+Inches(0.25), ly+Inches(0.2), cw-Inches(0.5), Inches(0.38),
             "업체 점검결과", 11, bold=True, color=RED)
    _textbox(sl, lx+Inches(0.25), ly+Inches(0.65), cw-Inches(0.5), Inches(0.65),
             "적합", 34, bold=True, color=INK)
    _textbox(sl, lx+Inches(0.25), ly+Inches(1.45), cw-Inches(0.5), Inches(2.0),
             "풍량 측정값에 대한 설계값 대비 정량 비교 없이 '적합'으로 판정.",
             16, color=MUTED, wrap=True)

    # 우측 카드 (AI - 초록)
    rx = Inches(7.0)
    _rect(sl, rx, ly, cw, ch, fill=WHITE, line_color=GREEN, line_w=2)
    _textbox(sl, rx+Inches(0.25), ly+Inches(0.2), cw-Inches(0.5), Inches(0.38),
             "AI 검토의견", 11, bold=True, color=GREEN)
    _textbox(sl, rx+Inches(0.25), ly+Inches(0.65), cw-Inches(0.5), Inches(0.65),
             "보완필요", 34, bold=True, color=INK)
    _textbox(sl, rx+Inches(0.25), ly+Inches(1.45), cw-Inches(0.5), Inches(2.0),
             "측정 18,634 CMH = 설계 21,000 CMH의 88.73%로 기준(±10%, 90–110%) 불충족. 정량 재판단 필요.",
             16, color=MUTED, wrap=True)

    _footer(sl)


# ══════════════════════════════════════════════════════════════
# 슬라이드 8: 기대효과 (dark)
# ══════════════════════════════════════════════════════════════

def slide_impact(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, W, H, fill=NAVY)
    _rect(sl, 0, 0, W, Inches(0.07), fill=GREEN)
    _rect(sl, 0, Inches(7.35), W, Inches(0.15), fill=GREEN)

    _textbox(sl, Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.55),
             "기대효과", 16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    _textbox(sl, Inches(1.5), Inches(2.1), Inches(10.3), Inches(0.9),
             "검토 일관성 · 업무 효율 · 데이터 축적",
             38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _accent_rule(sl, Inches(5.5), Inches(3.15), Inches(2.33), Inches(0.06))

    _textbox(sl, Inches(1.5), Inches(3.4), Inches(10.3), Inches(2.0),
             "표준화된 기준으로 검토 품질을 균질화하고, 검토 업무를 효율화하며,\n"
             "건물 성능 데이터를 디지털로 축적·활용할 수 있습니다.",
             20, color=LIGHT_B, align=PP_ALIGN.CENTER, wrap=True)

    _textbox(sl, Inches(1.5), Inches(6.9), Inches(10.3), Inches(0.4),
             "감사합니다", 16, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# 폰트 내장 (pptx_maker.py 로직 재사용)
# ══════════════════════════════════════════════════════════════

def _obfuscate(data: bytes, guid: str) -> bytes:
    key = bytes.fromhex(guid.replace("{","").replace("}","").replace("-",""))
    d   = bytearray(data)
    for i in range(min(32, len(d))):
        d[i] ^= key[15 - (i % 16)]
    return bytes(d)


def embed_fonts(pptx_bytes: bytes) -> bytes:
    reg_path  = os.path.join(FONT_DIR, "Pretendard-Regular.ttf")
    bold_path = os.path.join(FONT_DIR, "Pretendard-Bold.ttf")
    if not (os.path.exists(reg_path) and os.path.exists(bold_path)):
        print("  [경고] Pretendard 폰트 파일 없음 → 폰트 내장 생략")
        return pptx_bytes

    guid_r = str(uuid.uuid4()).upper()
    guid_b = str(uuid.uuid4()).upper()
    with open(reg_path,  "rb") as f: reg_data  = _obfuscate(f.read(), guid_r)
    with open(bold_path, "rb") as f: bold_data = _obfuscate(f.read(), guid_b)

    RELS_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"
    CT_FONT   = "application/vnd.openxmlformats-officedocument.obfuscatedFont"

    src = io.BytesIO(pptx_bytes)
    dst = io.BytesIO()
    with zipfile.ZipFile(src,"r") as zin, zipfile.ZipFile(dst,"w",zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)

            if item.filename == "ppt/presentation.xml":
                root = etree.fromstring(data)
                pPr  = root.find(qn("p:presentationPr"))
                if pPr is None: pPr = etree.SubElement(root, qn("p:presentationPr"))
                lst = etree.SubElement(pPr, qn("p:embeddedFontLst"))
                ef  = etree.SubElement(lst, qn("p:embeddedFont"))
                fe  = etree.SubElement(ef,  qn("p:font"))
                fe.set("typeface", FONT); fe.set("charset", "0")
                for tag, rid in [("p:regular","rIdFontR"),("p:bold","rIdFontB")]:
                    el = etree.SubElement(ef, qn(tag))
                    el.set("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id", rid)
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

            elif item.filename == "ppt/_rels/presentation.xml.rels":
                root = etree.fromstring(data)
                for rid, fname in [("rIdFontR","font1.fntdata"),("rIdFontB","font2.fntdata")]:
                    rel = etree.SubElement(root, "Relationship")
                    rel.set("Id", rid); rel.set("Type", RELS_TYPE); rel.set("Target", f"fonts/{fname}")
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

            elif item.filename == "[Content_Types].xml":
                root = etree.fromstring(data)
                NS   = "http://schemas.openxmlformats.org/package/2006/content-types"
                seen = {el.get("PartName") for el in root.findall(f"{{{NS}}}Override")}
                for fname in ("font1.fntdata","font2.fntdata"):
                    pn = f"/ppt/fonts/{fname}"
                    if pn not in seen:
                        ov = etree.SubElement(root, f"{{{NS}}}Override")
                        ov.set("PartName", pn); ov.set("ContentType", CT_FONT)
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

            zout.writestr(item, data)

        zout.writestr("ppt/fonts/font1.fntdata", reg_data)
        zout.writestr("ppt/fonts/font2.fntdata", bold_data)

    dst.seek(0)
    return dst.read()


# ══════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("슬라이드 생성 중...")
    slide_cover(prs)      ; print("  1/8 표지")
    slide_background(prs) ; print("  2/8 배경 및 필요성")
    slide_purpose(prs)    ; print("  3/8 목적")
    slide_process(prs)    ; print("  4/8 시스템 프로세스")
    slide_principles(prs) ; print("  5/8 검토 원칙")
    slide_results(prs)    ; print("  6/8 결과·실증")
    slide_case(prs)       ; print("  7/8 대표 검토 사례")
    slide_impact(prs)     ; print("  8/8 기대효과")

    buf = io.BytesIO()
    prs.save(buf)
    result = embed_fonts(buf.getvalue())

    out = os.path.join(os.path.dirname(__file__), "기계설비-성능점검-검토자문시스템.pptx")
    with open(out, "wb") as f:
        f.write(result)
    print(f"\n완료: {out}")


if __name__ == "__main__":
    build()
